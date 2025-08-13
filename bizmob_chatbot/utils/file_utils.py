"""
파일 처리 유틸리티
다양한 문서 형식을 처리하고 텍스트를 추출하는 기능
"""

import os
import io
import tempfile
from typing import List, Dict, Any, Optional
from langchain_core.documents import Document
import warnings
import streamlit as st

# 경고 억제
warnings.filterwarnings("ignore")


def get_supported_extensions() -> List[str]:
    """지원되는 파일 확장자 목록 반환"""
    return ['pdf', 'txt', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt']


def save_uploaded_file(uploaded_file, folder_path: str = "PDF_bizMOB_Guide") -> Optional[str]:
    """업로드된 파일을 지정된 폴더에 저장하고 파일 경로를 반환"""
    try:
        # 폴더가 없으면 생성 시도
        if not os.path.exists(folder_path):
            try:
                os.makedirs(folder_path, exist_ok=True)
            except PermissionError:
                # 권한 오류 시 임시 디렉토리 사용
                folder_path = tempfile.gettempdir()
                st.warning(f"권한 문제로 인해 임시 디렉토리에 파일을 저장합니다: {folder_path}")
        
        # 파일 경로 생성
        file_path = os.path.join(folder_path, uploaded_file.name)
        
        # 파일 저장
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # 업로드된 폴더 경로를 session_state에 저장
        if 'uploaded_folders' not in st.session_state:
            st.session_state.uploaded_folders = set()
        st.session_state.uploaded_folders.add(folder_path)
        
        return file_path
    except Exception as e:
        st.error(f"파일 저장 중 오류: {str(e)}")
        return None


def load_all_documents_from_folder(folder_path: str = "PDF_bizMOB_Guide") -> List[Document]:
    """지정된 폴더에서 모든 문서를 로드"""
    documents = []
    
    if not os.path.exists(folder_path):
        return documents
    
    try:
        for filename in os.listdir(folder_path):
            file_path = os.path.join(folder_path, filename)
            
            if os.path.isfile(file_path):
                file_extension = filename.split('.')[-1].lower()
                
                if file_extension in get_supported_extensions():
                    try:
                        with open(file_path, 'rb') as f:
                            # 파일 객체 생성 (StreamlitUploadedFile와 유사한 인터페이스)
                            class MockUploadedFile:
                                def __init__(self, file_path, filename):
                                    self.file_path = file_path
                                    self.name = filename
                                
                                def read(self):
                                    with open(self.file_path, 'rb') as f:
                                        return f.read()
                                
                                def getbuffer(self):
                                    with open(self.file_path, 'rb') as f:
                                        return f.read()
                            
                            mock_file = MockUploadedFile(file_path, filename)
                            file_docs = process_file(mock_file)
                            documents.extend(file_docs)
                            
                    except Exception as e:
                        st.warning(f"파일 {filename} 처리 중 오류: {str(e)}")
                        continue
        
        return documents
        
    except Exception as e:
        st.error(f"폴더 {folder_path} 로드 중 오류: {str(e)}")
        return documents


def process_file(uploaded_file) -> List[Document]:
    """업로드된 파일을 처리하여 Document 객체 리스트 반환"""
    
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return process_pdf(uploaded_file)
    elif file_extension in ['docx', 'doc']:
        return process_word(uploaded_file)
    elif file_extension in ['xlsx', 'xls']:
        return process_excel(uploaded_file)
    elif file_extension in ['pptx', 'ppt']:
        return process_powerpoint(uploaded_file)
    elif file_extension == 'txt':
        return process_text(uploaded_file)
    else:
        raise ValueError(f"지원되지 않는 파일 형식: {file_extension}")


def process_pdf(uploaded_file) -> List[Document]:
    """PDF 파일 처리"""
    try:
        import PyPDF2
        
        # 파일 내용 읽기
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
        
        documents = []
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text.strip():
                doc = Document(
                    page_content=text,
                    metadata={
                        'source': uploaded_file.name,
                        'page': page_num + 1,
                        'file_type': 'pdf'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"PDF 파일 처리 실패: {e}")


def process_word(uploaded_file) -> List[Document]:
    """Word 파일 처리"""
    try:
        from docx import Document as DocxDocument
        
        # 파일 내용 읽기
        doc = DocxDocument(io.BytesIO(uploaded_file.read()))
        
        documents = []
        text_content = []
        
        # 모든 단락에서 텍스트 추출
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # 표에서 텍스트 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        # 전체 텍스트를 하나의 문서로 생성
        if text_content:
            full_text = '\n\n'.join(text_content)
            doc = Document(
                page_content=full_text,
                metadata={
                    'source': uploaded_file.name,
                    'file_type': 'word'
                }
            )
            documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"Word 파일 처리 실패: {e}")


def process_excel(uploaded_file) -> List[Document]:
    """Excel 파일 처리"""
    try:
        import pandas as pd
        
        # 파일 내용 읽기
        df = pd.read_excel(io.BytesIO(uploaded_file.read()))
        
        documents = []
        
        # 각 행을 텍스트로 변환
        for index, row in df.iterrows():
            row_text = ' | '.join([str(cell) for cell in row if pd.notna(cell)])
            if row_text.strip():
                doc = Document(
                    page_content=row_text,
                    metadata={
                        'source': uploaded_file.name,
                        'row': index + 1,
                        'file_type': 'excel'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"Excel 파일 처리 실패: {e}")


def process_powerpoint(uploaded_file) -> List[Document]:
    """PowerPoint 파일 처리"""
    try:
        from pptx import Presentation
        
        # 파일 내용 읽기
        prs = Presentation(io.BytesIO(uploaded_file.read()))
        
        documents = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # 슬라이드의 모든 도형에서 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                full_text = '\n'.join(slide_text)
                doc = Document(
                    page_content=full_text,
                    metadata={
                        'source': uploaded_file.name,
                        'slide': slide_num + 1,
                        'file_type': 'powerpoint'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"PowerPoint 파일 처리 실패: {e}")


def process_text(uploaded_file) -> List[Document]:
    """텍스트 파일 처리"""
    try:
        # 파일 내용 읽기
        content = uploaded_file.read().decode('utf-8')
        
        doc = Document(
            page_content=content,
            metadata={
                'source': uploaded_file.name,
                'file_type': 'text'
            }
        )
        
        return [doc]
        
    except Exception as e:
        raise Exception(f"텍스트 파일 처리 실패: {e}")


def split_text_into_chunks(text: str, max_chunk_size: int = 4000, overlap: int = 200) -> List[str]:
    """텍스트를 청크로 분할"""
    if len(text) <= max_chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + max_chunk_size
        
        # 문장 경계에서 분할
        if end < len(text):
            # 마지막 마침표나 줄바꿈을 찾아서 분할
            last_period = text.rfind('.', start, end)
            last_newline = text.rfind('\n', start, end)
            
            if last_period > start and last_period > last_newline:
                end = last_period + 1
            elif last_newline > start:
                end = last_newline + 1
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap
        if start >= len(text):
            break
    
    return chunks


def get_file_info(uploaded_file) -> Dict[str, Any]:
    """파일 정보 반환"""
    return {
        'name': uploaded_file.name,
        'size': len(uploaded_file.getbuffer()),
        'type': uploaded_file.type
    }


def validate_file(uploaded_file) -> bool:
    """파일 유효성 검사"""
    if not uploaded_file:
        return False
    
    file_extension = uploaded_file.name.split('.')[-1].lower()
    return file_extension in get_supported_extensions()


def process_multiple_files(uploaded_files) -> List[Document]:
    """여러 파일 처리"""
    all_documents = []
    
    for uploaded_file in uploaded_files:
        if validate_file(uploaded_file):
            try:
                documents = process_file(uploaded_file)
                all_documents.extend(documents)
            except Exception as e:
                st.warning(f"파일 {uploaded_file.name} 처리 중 오류: {str(e)}")
                continue
    
    return all_documents 