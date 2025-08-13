#!/usr/bin/env python3
"""
bizMOB 챗봇 - ChromaDB 전용 버전
bizmob_chatbot_original.py의 모든 기능을 ChromaDB를 사용하여 구현
"""

import streamlit as st
from streamlit.runtime.uploaded_file_manager import UploadedFile
import os
import sys
import json
import pandas as pd
import re
import logging
import glob
import shutil
import subprocess
import time
import warnings
from datetime import datetime
from typing import List, Dict, Any, Optional
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document as DocxDocument
from safetensors.torch import load_file
import sentence_transformers
import html

# 경고 억제
warnings.filterwarnings("ignore")

# 환경 변수 설정 - safetensors 강제 사용
os.environ['TORCH_WARN_ON_LOAD'] = '0'
os.environ['TORCH_LOAD_WARN_ONLY'] = '0'
os.environ['PYTORCH_DISABLE_WARNINGS'] = '1'
os.environ["TOKENIZERS_PARALLELISM"] = "false"
os.environ['SAFETENSORS_FAST_GPU'] = '1'
os.environ['TRANSFORMERS_OFFLINE'] = '0'
os.environ['TORCH_WEIGHTS_ONLY'] = '1'
os.environ['HF_HUB_DISABLE_TELEMETRY'] = '1'
os.environ['TRANSFORMERS_USE_SAFETENSORS'] = '1'
# 추가 safetensors 강제 설정
os.environ['SAFETENSORS_FAST_GPU'] = '1'
os.environ['TRANSFORMERS_SAFE_SERIALIZATION'] = '1'
os.environ['HF_HUB_ENABLE_HF_TRANSFER'] = '1'
os.environ['TRANSFORMERS_CACHE'] = './model_cache'
os.environ['HF_HOME'] = './huggingface'
os.environ['TORCH_HOME'] = './torch_cache'

# 외부 소스 확장자 리스트 상수 선언
EXTERNAL_SOURCE_EXTS = [
    ".py", ".js", ".scss", ".ts", ".vue", ".md", ".txt", ".rst", ".json", ".yaml", ".yml"
]

# 전역 하이브리드 검색 설정 변수
HYBRID_SEARCH_CONFIG = {
    'bm25_weight': 0.3,
    'vector_weight': 0.7,
    'initial_k': 8,
    'final_k': 3,
    'enable_reranking': True,
    'metadata_boost': True,
    'recency_boost': True
}

# ChromaDB 사용 가능 여부 확인
try:
    import chromadb
    from langchain_community.vectorstores import Chroma
    CHROMADB_AVAILABLE = True
except ImportError:
    CHROMADB_AVAILABLE = False
    print("ChromaDB가 설치되지 않았습니다. pip install chromadb를 실행해주세요.")

# 로깅 설정
def setup_logging():
    """로깅 설정"""
    # 절대 경로로 로그 디렉토리 설정
    current_dir = os.path.dirname(os.path.abspath(__file__))
    log_dir = os.path.join(current_dir, "logs")
    
    # 로거 설정
    logger = logging.getLogger(__name__)
    logger.setLevel(logging.INFO)
    
    # 기존 핸들러 제거
    for handler in logger.handlers[:]:
        logger.removeHandler(handler)
    
    try:
        # 로그 디렉토리 생성 시도
        os.makedirs(log_dir, exist_ok=True)
        log_file = os.path.join(log_dir, f"bizmob_chatbot_{datetime.now().strftime('%Y%m%d')}.log")
        
        # 파일 핸들러 (UTF-8 인코딩)
        file_handler = logging.FileHandler(log_file, encoding='utf-8', mode='a')
        file_handler.setLevel(logging.INFO)
        file_formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(levelname)s - %(message)s')
        file_handler.setFormatter(file_formatter)
        logger.addHandler(file_handler)
        
        # 로그 파일 경로 출력
        print(f"로그 파일 경로: {log_file}")
        
    except (PermissionError, OSError) as e:
        # 로그 디렉토리 생성 실패 시 콘솔 로깅만 사용
        print(f"로그 디렉토리 생성 실패: {e}. 콘솔 로깅만 사용합니다.")
    
    # 콘솔 핸들러 (항상 추가)
    console_handler = logging.StreamHandler()
    console_handler.setLevel(logging.INFO)
    console_formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
    console_handler.setFormatter(console_formatter)
    logger.addHandler(console_handler)
    
    return logger

# 로거 초기화
logger = setup_logging()

# UI 컴포넌트 모듈 import
from ui_components import (
    apply_css_styles, setup_page_config,
    show_model_selector,
    show_embedding_model_info,
    show_admin_interface
)

# PyTorch 스레드 설정 (성능 최적화)
try:
    import torch
    torch.set_num_threads(1)
except Exception as e:
    logger.warning(f"PyTorch thread setup failed: {e}")

# 기타 필요한 import들
try:
    from langchain_core.documents.base import Document
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain.prompts import PromptTemplate
    from langchain_core.runnables import Runnable, RunnablePassthrough
    from langchain.schema.output_parser import StrOutputParser
    from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredExcelLoader, UnstructuredPowerPointLoader, UnstructuredWordDocumentLoader
    from langchain_huggingface import HuggingFaceEmbeddings
    from langchain_ollama import OllamaLLM
    from langchain_community.llms import Ollama
    from langchain.chains import RetrievalQA
    from langchain.retrievers import ParentDocumentRetriever
    from langchain.storage import InMemoryStore
    from langchain_core.embeddings import Embeddings
    from langchain_community.retrievers import BM25Retriever
    from langchain.retrievers import EnsembleRetriever
    from langchain.retrievers import ContextualCompressionRetriever
    from langchain.retrievers.document_compressors import LLMChainExtractor
except ImportError as e:
    st.error(f"필요한 라이브러리가 설치되지 않았습니다: {e}")
    st.stop()

# 환경변수 불러오기
from dotenv import load_dotenv, dotenv_values
load_dotenv()

# UI 스타일 및 페이지 설정 적용
apply_css_styles()
setup_page_config()



############################### 1단계 : 파일 업로드 및 관리 함수들 ##########################

def save_uploaded_file(uploaded_file: UploadedFile, folder_path: str = "PDF_bizMOB_Guide") -> str:
    """업로드된 파일을 지정된 폴더에 저장하고 파일 경로를 반환"""
    try:
        # 폴더가 없으면 생성
        if not os.path.exists(folder_path):
            os.makedirs(folder_path)
        
        # 파일 경로 생성
        file_path = os.path.join(folder_path, uploaded_file.name)
        
        # 파일 저장
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        return file_path
    except Exception as e:
        st.error(f"파일 저장 중 오류: {str(e)}")
        return None

def get_supported_file_types() -> dict:
    """지원하는 파일 형식과 설명을 반환"""
    return {
        'pdf': 'PDF 문서 (.pdf)',
        'xlsx': 'Excel 스프레드시트 (.xlsx)',
        'xls': 'Excel 스프레드시트 (.xls)',
        'pptx': 'PowerPoint 프레젠테이션 (.pptx)',
        'ppt': 'PowerPoint 프레젠테이션 (.ppt)',
        'docx': 'Word 문서 (.docx)',
        'doc': 'Word 문서 (.doc)'
    }

def validate_file_type(filename: str) -> bool:
    """파일 형식이 지원되는지 확인"""
    supported_extensions = ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.docx', '.doc']
    file_ext = os.path.splitext(filename.lower())[1]
    return file_ext in supported_extensions

def upload_and_process_files() -> bool:
    """파일 업로드 및 처리 함수"""
    st.markdown("### 📁 파일 업로드")
    st.markdown("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
    
    # 파일 업로드 위젯
    uploaded_files = st.file_uploader(
        "문서 파일을 선택하세요",
        type=['pdf', 'xlsx', 'xls', 'pptx', 'ppt', 'docx', 'doc'],
        accept_multiple_files=True,
        help="여러 파일을 동시에 선택할 수 있습니다."
    )
    
    if uploaded_files:
        st.markdown("#### 📋 업로드된 파일 목록")
        
        success_count = 0
        error_count = 0
        
        for uploaded_file in uploaded_files:
            col1, col2, col3 = st.columns([3, 1, 1])
            
            with col1:
                st.write(f"📄 {uploaded_file.name}")
            
            with col2:
                file_size = len(uploaded_file.getbuffer()) / 1024  # KB
                st.write(f"{file_size:.1f} KB")
            
            with col3:
                if validate_file_type(uploaded_file.name):
                    # 파일 저장
                    saved_path = save_uploaded_file(uploaded_file)
                    if saved_path:
                        st.success("✅")
                        success_count += 1
                    else:
                        st.error("❌")
                        error_count += 1
                else:
                    st.error("❌")
                    error_count += 1
        
        # 결과 요약
        if success_count > 0:
            st.success(f"✅ {success_count}개 파일이 성공적으로 업로드되었습니다.")
            
            # 벡터 데이터베이스 재초기화 옵션
            if st.button("🔄 업로드된 파일로 벡터 데이터베이스 재초기화", type="primary"):
                if initialize_vector_db_with_documents():
                    st.session_state.vector_db_initialized = True
                    st.success("벡터 데이터베이스가 새로운 파일들로 성공적으로 재초기화되었습니다!")
                    st.rerun()
                else:
                    st.error("벡터 데이터베이스 초기화에 실패했습니다.")
        
        if error_count > 0:
            st.error(f"❌ {error_count}개 파일 업로드에 실패했습니다.")
        
        return success_count > 0
    
    return False

def list_uploaded_files(folder_path: str = "PDF_bizMOB_Guide") -> dict:
    """업로드된 파일들을 형식별로 분류하여 반환"""
    if not os.path.exists(folder_path):
        return {}
    
    files_by_type = {
        'PDF': [],
        'Excel': [],
        'PowerPoint': [],
        'Word': []
    }
    
    # 지원하는 파일 확장자들
    supported_extensions = {
        '*.pdf': 'PDF',
        '*.xlsx': 'Excel',
        '*.xls': 'Excel',
        '*.pptx': 'PowerPoint',
        '*.ppt': 'PowerPoint',
        '*.docx': 'Word',
        '*.doc': 'Word'
    }
    
    for pattern, file_type in supported_extensions.items():
        files = glob.glob(os.path.join(folder_path, pattern))
        for file_path in files:
            file_size = os.path.getsize(file_path) / 1024  # KB
            files_by_type[file_type].append({
                'name': os.path.basename(file_path),
                'path': file_path,
                'size': file_size
            })
    
    return files_by_type

def safe_key(filename: str) -> str:
    """파일명을 안전한 session state 키로 변환"""
    # 특수문자를 언더스코어로 변환
    safe_name = re.sub(r'[^a-zA-Z0-9_]', '_', filename)
    # 숫자로 시작하지 않도록 prefix 추가
    if safe_name and safe_name[0].isdigit():
        safe_name = 'file_' + safe_name
    return safe_name

def delete_file(file_path: str) -> bool:
    """파일 삭제 함수"""
    try:
        # 파일이 존재하는지 확인
        if not os.path.exists(file_path):
            st.error(f"파일이 존재하지 않습니다: {file_path}")
            return False
        
        # 파일 삭제
        os.remove(file_path)
        
        # 삭제 확인
        if os.path.exists(file_path):
            st.error(f"파일 삭제에 실패했습니다: {file_path}")
            return False
        
        return True
    except PermissionError:
        st.error(f"파일 삭제 권한이 없습니다: {file_path}")
        return False
    except Exception as e:
        st.error(f"파일 삭제 중 오류: {str(e)}")
        return False

def delete_file_with_confirmation(file_path: str, file_name: str) -> bool:
    """확인 후 파일 삭제 함수"""
    # 삭제 확인을 위한 session state 키
    confirm_key = f"confirm_delete_{file_name}"
    
    if confirm_key not in st.session_state:
        st.session_state[confirm_key] = False
    
    if not st.session_state[confirm_key]:
        # 삭제 확인 버튼
        if st.button(f"🗑️ 삭제 확인", key=f"confirm_{file_name}"):
            st.session_state[confirm_key] = True
            return False
        return False
    else:
        # 실제 삭제 실행
        if delete_file(file_path):
            st.success(f"✅ {file_name} 파일이 삭제되었습니다.")
            # session state 정리
            del st.session_state[confirm_key]
            # 벡터 데이터베이스 재초기화 제안
            st.warning("⚠️ 삭제된 파일이 벡터 데이터베이스에 반영되려면 재초기화가 필요합니다.")
            if st.button("🔄 벡터DB 재초기화", key=f"reinit_after_delete_{file_name}"):
                if initialize_vector_db():
                    st.session_state.vector_db_initialized = True
                    st.success("벡터 데이터베이스가 성공적으로 재초기화되었습니다!")
            return True
        else:
            st.error(f"❌ {file_name} 파일 삭제에 실패했습니다.")
            # session state 정리
            del st.session_state[confirm_key]
            return False

def manage_uploaded_files() -> None:
    """업로드된 파일 관리 인터페이스"""
    st.markdown("### 📂 업로드된 파일 관리")
    
    # 신규 파일 업로드 섹션 추가
    with st.expander("📁 신규 파일 업로드", expanded=False):
        st.markdown("#### 📤 파일 업로드")
        st.markdown("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
        
        # 파일 업로드 위젯
        uploaded_files = st.file_uploader(
            "문서 파일을 선택하세요",
            type=['pdf', 'xlsx', 'xls', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="여러 파일을 동시에 선택할 수 있습니다.",
            key="file_manager_uploader"
        )
        
        if uploaded_files:
            st.markdown("#### 📋 업로드할 파일 목록")
            
            success_count = 0
            error_count = 0
            
            for uploaded_file in uploaded_files:
                col1, col2, col3 = st.columns([3, 1, 1])
                
                with col1:
                    st.write(f"📄 {uploaded_file.name}")
                
                with col2:
                    file_size = len(uploaded_file.getbuffer()) / 1024  # KB
                    st.write(f"{file_size:.1f} KB")
                
                with col3:
                    if validate_file_type(uploaded_file.name):
                        # 파일 저장
                        saved_path = save_uploaded_file(uploaded_file)
                        if saved_path:
                            st.success("✅")
                            success_count += 1
                        else:
                            st.error("❌")
                            error_count += 1
                    else:
                        st.error("❌")
                        error_count += 1
            
            # 결과 요약
            if success_count > 0:
                st.success(f"✅ {success_count}개 파일이 성공적으로 업로드되었습니다.")
                
                # 벡터 데이터베이스 재초기화 옵션
                col1, col2 = st.columns([1, 1])
                with col1:
                    if st.button("🔄 벡터DB 재초기화", type="primary", key="file_manager_reinit"):
                        if initialize_vector_db_with_documents():
                            st.session_state.vector_db_initialized = True
                            st.success("벡터 데이터베이스가 성공적으로 재초기화되었습니다!")
                        else:
                            st.error("벡터 데이터베이스 초기화에 실패했습니다.")
                
                with col2:
                    if st.button("🔄 페이지 새로고침", key="file_manager_refresh"):
                        pass  # 페이지 새로고침은 브라우저에서 직접 수행
            
            if error_count > 0:
                st.error(f"❌ {error_count}개 파일 업로드에 실패했습니다.")
    
    st.markdown("---")
    
    files_by_type = list_uploaded_files()
    
    if not any(files_by_type.values()):
        st.info("업로드된 파일이 없습니다.")
        st.markdown("---")
        st.markdown("### 📋 지원 파일 형식")
        st.markdown("- **PDF** (.pdf): 매뉴얼, 가이드 문서")
        st.markdown("- **Excel** (.xlsx, .xls): 데이터 시트, 분석 자료")
        st.markdown("- **PowerPoint** (.pptx, .ppt): 프레젠테이션, 교육 자료")
        st.markdown("- **Word** (.docx, .doc): 보고서, 문서, 매뉴얼")
        return
    
    # 파일 통계 정보
    total_files = sum(len(files) for files in files_by_type.values())
    total_size = sum(sum(file_info['size'] for file_info in files) for files in files_by_type.values())
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("📁 총 파일 수", f"{total_files}개")
    with col2:
        st.metric("💾 총 용량", f"{total_size:.1f} KB")
    with col3:
        if check_vector_db_exists():
            st.success("✅ 벡터DB 준비됨")
        else:
            st.warning("⚠️ 벡터DB 미준비")
    
    # 벡터DB 재초기화 버튼
    if st.button("🔄 벡터 데이터베이스 재초기화", type="primary", key="file_manager_main_reinit"):
        if initialize_vector_db_with_documents():
            st.session_state.vector_db_initialized = True
            st.success("벡터 데이터베이스가 성공적으로 재초기화되었습니다!")
        else:
            st.error("벡터 데이터베이스 초기화에 실패했습니다.")
    
    st.markdown("---")
    
    # 파일 형식별로 탭 생성 (수정)
    file_types_with_files = [(file_type, files) for file_type, files in files_by_type.items() if files]
    tab_names = [f"{file_type} ({len(files)})" for file_type, files in file_types_with_files]
    if tab_names:
        tabs = st.tabs(tab_names)
        for i, (file_type, files) in enumerate(file_types_with_files):
            with tabs[i]:
                st.markdown(f"#### {file_type} 파일 목록")
                
                for file_info in files:
                    with st.expander(f"📄 {file_info['name']} ({file_info['size']:.1f} KB)"):
                        col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                        
                        with col1:
                            st.write(f"**파일명:** {file_info['name']}")
                            st.write(f"**크기:** {file_info['size']:.1f} KB")
                            st.write(f"**경로:** {file_info['path']}")
                        
                        with col2:
                            # 파일 다운로드 버튼
                            safe_download_key = safe_key(file_info['name'])
                            with open(file_info['path'], 'rb') as f:
                                st.download_button(
                                    label="📥 다운로드",
                                    data=f.read(),
                                    file_name=file_info['name'],
                                    key=f"download_{safe_download_key}"
                                )
                        
                        with col3:
                            # 파일 미리보기 버튼
                            safe_preview_key = safe_key(file_info['name'])
                            if st.button("👁️ 미리보기", key=f"preview_{safe_preview_key}"):
                                st.session_state.preview_file = file_info['path']
                                st.session_state.preview_file_type = file_type
                                st.session_state.preview_file_name = file_info['name']
                        
                        with col4:
                            # 파일 삭제 버튼 - 안전한 키 사용
                            safe_file_key = safe_key(file_info['name'])
                            delete_key = f"delete_{safe_file_key}"
                            
                            if delete_key not in st.session_state:
                                st.session_state[delete_key] = False
                            
                            if not st.session_state[delete_key]:
                                if st.button("🗑️ 삭제", key=f"btn_{safe_file_key}"):
                                    st.session_state[delete_key] = True
                            else:
                                st.warning(f"정말로 '{file_info['name']}' 파일을 삭제하시겠습니까?")
                                col_confirm1, col_confirm2 = st.columns(2)
                                with col_confirm1:
                                    if st.button("✅ 확인", key=f"confirm_{safe_file_key}"):
                                        if delete_file(file_info['path']):
                                            st.success(f"✅ {file_info['name']} 파일이 삭제되었습니다.")
                                            # 벡터DB 재초기화 제안
                                            st.warning("⚠️ 삭제된 파일이 벡터 데이터베이스에 반영되려면 재초기화가 필요합니다.")
                                            if st.button("🔄 벡터DB 재초기화", key=f"reinit_{safe_file_key}"):
                                                if initialize_vector_db_with_documents():
                                                    st.session_state.vector_db_initialized = True
                                                    st.success("벡터 데이터베이스가 성공적으로 재초기화되었습니다!")
                                            # session state 정리
                                            del st.session_state[delete_key]
                                        else:
                                            st.error(f"❌ {file_info['name']} 파일 삭제에 실패했습니다.")
                                            del st.session_state[delete_key]
                                with col_confirm2:
                                    if st.button("❌ 취소", key=f"cancel_{safe_file_key}"):
                                        del st.session_state[delete_key]
                    
                    # 파일 미리보기 섹션
                    if 'preview_file' in st.session_state and st.session_state.preview_file_type == file_type:
                        st.markdown("---")
                        st.markdown(f"#### 👁️ 파일 미리보기: {st.session_state.preview_file_name}")
                        
                        try:
                            if file_type == 'PDF':
                                # PDF 미리보기
                                images = convert_pdf_to_images(st.session_state.preview_file)
                                if images:
                                    st.image(images[0], caption="첫 페이지 미리보기", width=400)
                                    st.info(f"총 {len(images)}페이지")
                            
                            elif file_type == 'Excel':
                                # Excel 미리보기
                                excel_file = pd.ExcelFile(st.session_state.preview_file)
                                sheet_names = excel_file.sheet_names
                                st.write(f"**시트 목록:** {', '.join(sheet_names)}")
                                
                                if sheet_names:
                                    selected_sheet = st.selectbox("시트 선택", sheet_names)
                                    df = pd.read_excel(st.session_state.preview_file, sheet_name=selected_sheet)
                                    st.dataframe(df.head(10), use_container_width=True)
                                    st.info(f"총 {len(df)}행, {len(df.columns)}열")
                            
                            elif file_type == 'PowerPoint':
                                # PowerPoint 미리보기
                                prs = Presentation(st.session_state.preview_file)
                                st.write(f"**총 슬라이드 수:** {len(prs.slides)}")
                                
                                if prs.slides:
                                    slide_content = ""
                                    for i, slide in enumerate(prs.slides[:3]):  # 처음 3개 슬라이드만
                                        slide_content += f"**슬라이드 {i+1}:**\n"
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text") and shape.text.strip():
                                                slide_content += f"{shape.text}\n"
                                        slide_content += "\n"
                                    
                                    st.text_area("슬라이드 내용 미리보기", slide_content, height=200)
                            
                            elif file_type == 'Word':
                                # Word 미리보기
                                doc = DocxDocument(st.session_state.preview_file)
                                st.write(f"**제목:** {doc.core_properties.title or '제목 없음'}")
                                st.write(f"**작성자:** {doc.core_properties.author or '작성자 없음'}")
                                
                                # 처음 몇 단락만 미리보기
                                preview_text = ""
                                for para in doc.paragraphs[:10]:
                                    if para.text.strip():
                                        preview_text += para.text + "\n\n"
                                
                                st.text_area("문서 내용 미리보기", preview_text, height=200)
                        
                        except Exception as e:
                            st.error(f"파일 미리보기 중 오류: {str(e)}")
                        
                        # 미리보기 닫기 버튼
                        if st.button("❌ 미리보기 닫기"):
                            del st.session_state.preview_file
                            del st.session_state.preview_file_type
                            del st.session_state.preview_file_name 

############################### 1단계 : PDF 문서를 벡터DB에 저장하는 함수들 ##########################

## 1: PDF_bizMOB_Guide 폴더에서 모든 문서 파일을 찾아서 Document로 변환
def load_all_documents_from_folder(folder_path: str = "PDF_bizMOB_Guide") -> List[Document]:
    documents = []
    
    if not os.path.exists(folder_path):
        st.error(f"폴더가 존재하지 않습니다: {folder_path}")
        return documents
    
    # 지원하는 파일 확장자들
    supported_extensions = {
        '*.pdf': 'PDF',
        '*.xlsx': 'Excel',
        '*.xls': 'Excel',
        '*.pptx': 'PowerPoint',
        '*.ppt': 'PowerPoint',
        '*.docx': 'Word',
        '*.doc': 'Word'
    }
    
    all_files = []
    for pattern, file_type in supported_extensions.items():
        files = glob.glob(os.path.join(folder_path, pattern))
        all_files.extend([(f, file_type) for f in files])
    
    if not all_files:
        st.warning(f"{folder_path} 폴더에 지원하는 문서 파일이 없습니다.")
        st.info("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
        return documents
    
    for file_path, file_type in all_files:
        try:
            st.info(f"{file_type} 파일 로딩 중: {os.path.basename(file_path)}")
            
            if file_type == 'PDF':
                loader = PyMuPDFLoader(file_path)
                doc = loader.load()
            elif file_type == 'Excel':
                doc = load_excel_file(file_path)
            elif file_type == 'PowerPoint':
                doc = load_powerpoint_file(file_path)
            elif file_type == 'Word':
                doc = load_word_file(file_path)
            
            for d in doc:
                d.metadata['file_path'] = file_path
                d.metadata['file_name'] = os.path.basename(file_path)
                d.metadata['file_type'] = file_type
            
            documents.extend(doc)
            st.success(f"✅ {os.path.basename(file_path)} ({file_type}) 로딩 완료")
        except Exception as e:
            st.error(f"❌ {os.path.basename(file_path)} ({file_type}) 로딩 실패: {str(e)}")
    
    # 외부 소스 코드(.py, .md 등)도 포함
    ext_src_dir = "external_sources"
    if os.path.exists(ext_src_dir):
        for repo in os.listdir(ext_src_dir):
            repo_path = os.path.join(ext_src_dir, repo)
            for ext in EXTERNAL_SOURCE_EXTS:
                for file in glob.glob(f"{repo_path}/**/*{ext}", recursive=True):
                    try:
                        if os.path.getsize(file) > 2*1024*1024:
                            st.warning(f"{file} 파일은 2MB를 초과하여 벡터DB에 포함되지 않습니다.")
                            continue
                        with open(file, "r", encoding="utf-8", errors="ignore") as f:
                            content = f.read()
                        doc = Document(
                            page_content=content,
                            metadata={
                                'file_path': file,
                                'file_name': os.path.basename(file),
                                'file_type': 'Source',
                                'repo': repo
                            }
                        )
                        documents.append(doc)
                    except Exception as e:
                        st.warning(f"{file} 파일을 읽는 중 오류 발생: {e}")
    return documents

def load_excel_file(file_path: str) -> List[Document]:
    """Excel 파일을 로드하여 Document 리스트로 변환"""
    documents = []
    try:
        # Excel 파일 읽기
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            
            # 데이터프레임을 텍스트로 변환
            text_content = f"시트명: {sheet_name}\n\n"
            
            # 헤더 정보 추가
            if not df.empty:
                text_content += f"컬럼: {', '.join(df.columns.tolist())}\n\n"
                
                # 데이터 내용 추가 (처음 100행까지만)
                for idx, row in df.head(100).iterrows():
                    row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    if row_text.strip():
                        text_content += f"행 {idx+1}: {row_text}\n"
            
            # Document 생성
            doc = Document(
                page_content=text_content,
                metadata={
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_type': 'Excel',
                    'sheet_name': sheet_name
                }
            )
            documents.append(doc)
            
    except Exception as e:
        st.error(f"Excel 파일 처리 중 오류: {str(e)}")
    
    return documents

def load_powerpoint_file(file_path: str) -> List[Document]:
    """PowerPoint 파일을 로드하여 Document 리스트로 변환"""
    documents = []
    try:
        # PowerPoint 파일 읽기
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content = f"슬라이드 {slide_num}:\n\n"
            
            # 슬라이드의 모든 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += f"{shape.text}\n\n"
            
            # 텍스트가 있는 경우에만 Document 생성
            if text_content.strip() and text_content != f"슬라이드 {slide_num}:\n\n":
                doc = Document(
                    page_content=text_content,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_type': 'PowerPoint',
                        'slide_number': slide_num
                    }
                )
                documents.append(doc)
                
    except Exception as e:
        st.error(f"PowerPoint 파일 처리 중 오류: {str(e)}")
    
    return documents

def load_word_file(file_path: str) -> List[Document]:
    """Word 파일을 로드하여 Document 리스트로 변환"""
    documents = []
    try:
        # Word 파일 읽기
        doc = DocxDocument(file_path)
        
        # 전체 문서를 하나의 Document로 처리
        text_content = ""
        
        # 제목 정보 추가
        if doc.core_properties.title:
            text_content += f"제목: {doc.core_properties.title}\n\n"
        
        # 단락별로 텍스트 추출
        for para in doc.paragraphs:
            if para.text.strip():
                text_content += para.text + "\n\n"
        
        # 표 내용 추출
        for table in doc.tables:
            text_content += "표 내용:\n"
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                if row_text.strip():
                    text_content += row_text + "\n"
            text_content += "\n"
        
        # 텍스트가 있는 경우에만 Document 생성
        if text_content.strip():
            doc_chunk = Document(
                page_content=text_content,
                metadata={
                    'file_path': file_path,
                    'file_name': os.path.basename(file_path),
                    'file_type': 'Word',
                    'title': doc.core_properties.title or 'Unknown',
                    'author': doc.core_properties.author or 'Unknown'
                }
            )
            documents.append(doc_chunk)
            
    except Exception as e:
        st.error(f"Word 파일 처리 중 오류: {str(e)}")
    
    return documents

## 2: Document를 더 작은 document로 변환
def chunk_documents(documents: List[Document]) -> List[Document]:
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=100)
    return text_splitter.split_documents(documents)

## 3: Document를 벡터DB로 저장 (ChromaDB 사용)
def save_to_vector_store(documents: List[Document]) -> None:
    try:
        # 선택된 임베딩 모델 사용
        embeddings = get_embedding_model()
        selected_embedding = st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask')
        
        st.info(f"임베딩 모델 로딩 중: {selected_embedding}")
        
        # ChromaDB에 저장
        save_to_chroma_store(documents)
        st.success("✅ 벡터 데이터베이스 저장 완료 (ChromaDB 사용)")
    except Exception as e:
        st.error(f"❌ 벡터 데이터베이스 저장 실패: {str(e)}")

## 4: 벡터DB 초기화 함수 (문서 로딩 포함)
def initialize_vector_db_with_documents():
    """PDF_bizMOB_Guide 폴더의 모든 문서를 로드하여 벡터DB를 초기화"""
    with st.spinner("bizMOB Platform 가이드 문서들을 로딩하고 벡터 데이터베이스를 생성하는 중..."):
        # 모든 문서들 로드 (PDF, Excel, PowerPoint)
        documents = load_all_documents_from_folder()
        
        if not documents:
            st.error("로드할 문서가 없습니다.")
            return False
        
        # 문서 청킹
        st.info("문서를 청크로 분할하는 중...")
        chunked_documents = chunk_documents(documents)
        st.success(f"✅ {len(chunked_documents)}개의 청크로 분할 완료")
        
        # 벡터DB 저장
        st.info("벡터 데이터베이스에 저장하는 중...")
        save_to_chroma_store(chunked_documents)
        
        # 성공적으로 초기화된 모델 정보를 파일에 저장
        try:
            model_info = {
                'ai_model': st.session_state.get('selected_model', 'exaone3.5'),
                'embedding_model': st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask'),
                'timestamp': pd.Timestamp.now().isoformat()
            }
            
            import json
            with open('vector_db_model_info.json', 'w', encoding='utf-8') as f:
                json.dump(model_info, f, ensure_ascii=False, indent=2)
            
            st.success("✅ 모델 정보가 저장되었습니다.")
        except Exception as e:
            st.warning(f"⚠️ 모델 정보 저장 중 오류: {str(e)}")
        
        return True

## 4-1: 벡터DB 초기화 함수 (ChromaDB만 초기화)
def initialize_vector_db():
    """벡터 데이터베이스 초기화"""
    logger.info("Vector database initialization started")
    
    if not CHROMADB_AVAILABLE:
        error_msg = "ChromaDB가 설치되지 않았습니다."
        logger.error(error_msg)
        st.error(error_msg)
        return False
    
    try:
        # ChromaDB 디렉토리 생성
        chroma_path = get_chroma_db_path()
        logger.info(f"ChromaDB path: {chroma_path}")
        os.makedirs(chroma_path, exist_ok=True)
        logger.info("ChromaDB directory created successfully")
        
        # ChromaDB 클라이언트 생성 및 컬렉션 초기화
        import chromadb
        import time
        
        # 클라이언트 연결 시도
        max_retries = 3
        client = None
        
        for attempt in range(max_retries):
            try:
                client = chromadb.PersistentClient(
                    path=chroma_path,
                    settings=chromadb.config.Settings(
                        allow_reset=True,
                        anonymized_telemetry=False
                    )
                )
                break
            except Exception as e:
                if attempt < max_retries - 1:
                    logger.warning(f"Client connection attempt {attempt + 1} failed: {e}")
                    time.sleep(1)
                else:
                    raise e
        
        if client is None:
            raise Exception("ChromaDB 클라이언트 연결 실패")
        
        collection_name = "bizmob_documents"
        
        # 기존 컬렉션이 있으면 삭제하고 새로 생성
        try:
            client.delete_collection(name=collection_name)
            logger.info("Existing collection deleted")
            time.sleep(0.5)  # 잠시 대기
        except:
            logger.info("No existing collection to delete")
        
        # 새 컬렉션 생성
        collection = client.create_collection(name=collection_name)
        logger.info("New collection created")
        
        # 모델 정보 저장
        model_info = {
            'ai_model': st.session_state.get('selected_model', 'exaone3.5'),
            'embedding_model': st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask'),
            'timestamp': pd.Timestamp.now().isoformat()
        }
        logger.info(f"Model info: {model_info}")
        
        model_info_path = get_model_info_path()
        logger.info(f"Model info file path: {model_info_path}")
        
        with open(model_info_path, 'w', encoding='utf-8') as f:
            json.dump(model_info, f, ensure_ascii=False, indent=2)
        
        st.session_state.vector_db_initialized = True
        st.success("✅ ChromaDB 벡터 데이터베이스 초기화 완료")
        logger.info("Vector database initialization successful")
        return True
        
    except Exception as e:
        error_msg = f"벡터 데이터베이스 초기화 실패: {e}"
        logger.error(f"Vector database initialization failed: {e}", exc_info=True)
        st.error(f"❌ {error_msg}")
        return False

def save_to_chroma_store(documents: list) -> None:
    """문서를 ChromaDB에 저장"""
    logger.info(f"Vector database save started - document count: {len(documents)}")
    
    if not CHROMADB_AVAILABLE:
        error_msg = "ChromaDB가 설치되지 않았습니다."
        logger.error(error_msg)
        st.error(error_msg)
        return
    
    try:
        # 환경 변수 설정
        import os
        os.environ['TOKENIZERS_PARALLELISM'] = 'false'
        os.environ['PYTORCH_DISABLE_WARNINGS'] = '1'
        
        # 임베딩 모델 로드
        selected_embedding = st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask')
        logger.info(f"Loading embedding model: {selected_embedding}")
        
        try:
            # SafeSentenceTransformerEmbeddings 사용 (torch.load 취약점 방지)
            embeddings = SafeSentenceTransformerEmbeddings(
                model_name=selected_embedding,
                device='cpu'
            )
            logger.info("Embedding model loaded successfully")
        except Exception as e:
            logger.error(f"Embedding model loading failed: {e}")
            st.error(f"임베딩 모델 로딩 실패: {e}")
            st.info("다른 임베딩 모델을 시도해보세요.")
            return
        
        # ChromaDB 클라이언트 생성 (지속적 저장을 위해 경로 지정)
        try:
            import chromadb
            import time
            chroma_path = get_chroma_db_path()
            os.makedirs(chroma_path, exist_ok=True)
            
            # session_state에서 기존 클라이언트 확인
            if 'chroma_client' in st.session_state:
                try:
                    # 기존 클라이언트 해제 시도
                    del st.session_state.chroma_client
                    import gc
                    gc.collect()
                    time.sleep(0.5)
                except:
                    pass
            
            # 지속적 저장을 위한 ChromaDB 클라이언트 생성 (재시도 로직)
            max_retries = 3
            client = None
            
            for attempt in range(max_retries):
                try:
                    client = chromadb.PersistentClient(
                        path=chroma_path,
                        settings=chromadb.config.Settings(
                            allow_reset=True,
                            anonymized_telemetry=False
                        )
                    )
                    # session_state에 클라이언트 저장
                    st.session_state.chroma_client = client
                    break
                except Exception as e:
                    if attempt < max_retries - 1:
                        logger.warning(f"Client connection attempt {attempt + 1} failed: {e}")
                        time.sleep(1)
                    else:
                        raise e
            
            if client is None:
                raise Exception("ChromaDB 클라이언트 연결 실패")
            
            collection_name = "bizmob_documents"
            
            # 기존 컬렉션이 있으면 삭제하고 새로 생성
            try:
                client.delete_collection(name=collection_name)
                logger.info("Existing collection deleted")
                time.sleep(0.5)  # 잠시 대기
            except:
                logger.info("No existing collection to delete")
            
            # 새 컬렉션 생성
            collection = client.create_collection(name=collection_name)
            logger.info("New collection created")
            
            # 문서 텍스트와 메타데이터 추출 (텍스트 정제)
            documents_texts = []
            documents_metadatas = []
            documents_ids = []
            
            for i, doc in enumerate(documents):
                # 텍스트 정제 (특수문자 및 인코딩 문제 해결)
                clean_text = doc.page_content.strip()
                if clean_text:
                    documents_texts.append(clean_text)
                    documents_metadatas.append(doc.metadata)
                    documents_ids.append(f"doc_{i}_{datetime.now().strftime('%Y%m%d_%H%M%S')}")
            
            if not documents_texts:
                st.warning("저장할 문서가 없습니다.")
                return
            
            # 임베딩 생성
            logger.info("Generating embeddings...")
            embeddings_list = embeddings.embed_documents(documents_texts)
            logger.info(f"Generated {len(embeddings_list)} embeddings")
            
            # ChromaDB에 저장
            collection.add(
                documents=documents_texts,
                embeddings=embeddings_list,
                metadatas=documents_metadatas,
                ids=documents_ids
            )
            
            logger.info("ChromaDB document save completed")
            st.success(f"✅ {len(documents_texts)}개 문서가 벡터 데이터베이스에 저장되었습니다.")
            
            # 저장된 문서 수 확인
            try:
                count = collection.count()
                logger.info(f"Total documents in collection: {count}")
                st.info(f"📊 현재 벡터 DB에 저장된 문서 수: {count}개")
                
                # 저장된 문서 샘플 확인
                if count > 0:
                    sample_results = collection.get(limit=1)
                    if sample_results['documents']:
                        sample_text = sample_results['documents'][0][:100]
                        logger.info(f"Sample document: {sample_text}...")
                        st.info(f"📝 저장된 문서 샘플: {sample_text}...")
                
            except Exception as e:
                logger.warning(f"Could not get collection count: {e}")
            
        except Exception as e:
            logger.error(f"ChromaDB save failed: {e}")
            st.error(f"벡터 데이터베이스 저장 실패: {e}")
            return
        
    except Exception as e:
        error_msg = f"벡터 데이터베이스 저장 실패: {e}"
        logger.error(f"Vector database save failed: {e}", exc_info=True)
        st.error(f"❌ {error_msg}")

def load_chroma_store():
    """ChromaDB에서 벡터 스토어 로드 (강화된 단일 인스턴스 관리)"""
    if not CHROMADB_AVAILABLE:
        error_msg = "ChromaDB가 설치되지 않았습니다."
        logger.error(error_msg)
        st.error(error_msg)
        return None
    
    try:
        logger.info("=== ChromaDB 벡터 스토어 로드 시작 ===")
        
        # 캐시된 임베딩 모델 사용
        embeddings = get_embedding_model()
        if embeddings is None:
            logger.error("임베딩 모델 로드 실패")
            return None
        
        logger.info("임베딩 모델 로드 성공")
        
        # ChromaDB 클라이언트 생성 (강화된 단일 인스턴스 관리)
        import chromadb
        from langchain_community.vectorstores import Chroma
        import time
        import gc
        
        chroma_path = get_chroma_db_path()
        logger.info(f"ChromaDB 경로: {chroma_path}")
        
        # 디렉토리가 없으면 생성
        if not os.path.exists(chroma_path):
            os.makedirs(chroma_path, exist_ok=True)
            logger.info(f"ChromaDB 디렉토리 생성: {chroma_path}")
        
        # 현재 선택된 모델과 임베딩으로 고유 키 생성
        selected_model = st.session_state.get('selected_model', 'exaone3.5')
        selected_embedding = st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask')
        global_vector_store_key = f"global_vector_store_{selected_model}_{selected_embedding}"
        logger.info(f"벡터 스토어 키: {global_vector_store_key}")
        
        # 기존 벡터 스토어가 있으면 재사용
        if global_vector_store_key in st.session_state:
            try:
                logger.info("기존 벡터 스토어 캐시 확인 중...")
                vector_store = st.session_state[global_vector_store_key]
                # 간단한 테스트로 연결 상태 확인
                test_collection = vector_store._collection
                doc_count = test_collection.count()
                logger.info(f"기존 벡터 스토어 재사용 성공 - 문서 수: {doc_count}")
                return vector_store
            except Exception as e:
                logger.warning(f"기존 벡터 스토어 재사용 실패: {e}")
                # 기존 벡터 스토어 제거
                del st.session_state[global_vector_store_key]
                # 메모리 정리
                gc.collect()
                time.sleep(1)
        
        # ChromaDB 클라이언트 생성 (타임아웃 설정)
        logger.info("ChromaDB 클라이언트 생성 중...")
        try:
            client = chromadb.PersistentClient(
                path=chroma_path,
                settings=chromadb.config.Settings(
                    allow_reset=True,
                    anonymized_telemetry=False,
                    is_persistent=True,
                    persist_directory=chroma_path
                )
            )
            
            # 클라이언트 연결 테스트
            client.heartbeat()
            logger.info("ChromaDB 클라이언트 연결 성공")
            
        except Exception as e:
            logger.warning(f"ChromaDB 클라이언트 연결 실패, 재시도: {e}")
            # 기존 프로세스 정리 후 재시도
            try:
                import psutil
                logger.info("기존 ChromaDB 프로세스 정리 중...")
                for proc in psutil.process_iter(['pid', 'name', 'cmdline']):
                    try:
                        if 'chroma' in proc.info['name'].lower() or any('chroma' in str(cmd).lower() for cmd in proc.info['cmdline'] or []):
                            proc.terminate()
                            logger.info(f"ChromaDB 프로세스 종료: {proc.info['pid']}")
                    except (psutil.NoSuchProcess, psutil.AccessDenied):
                        pass
                time.sleep(2)
                
                logger.info("ChromaDB 클라이언트 재생성 중...")
                client = chromadb.PersistentClient(
                    path=chroma_path,
                    settings=chromadb.config.Settings(
                        allow_reset=True,
                        anonymized_telemetry=False,
                        is_persistent=True,
                        persist_directory=chroma_path
                    )
                )
                client.heartbeat()
                logger.info("ChromaDB 클라이언트 재연결 성공")
                
            except Exception as e2:
                logger.error(f"ChromaDB 클라이언트 재연결 실패: {e2}")
                return None
        
        collection_name = "bizmob_documents"
        logger.info(f"컬렉션 이름: {collection_name}")
        
        # 컬렉션 가져오기 (타임아웃 설정)
        try:
            collection = client.get_collection(name=collection_name)
            doc_count = collection.count()
            logger.info(f"기존 컬렉션 로드 성공 - 문서 수: {doc_count}")
        except Exception as e:
            logger.info("새 컬렉션 생성")
            collection = client.create_collection(name=collection_name)
            logger.info("새 컬렉션 생성 완료")
        
        # LangChain Chroma 벡터 스토어 생성
        logger.info("LangChain Chroma 벡터 스토어 생성 중...")
        vector_store = Chroma(
            client=client,
            collection_name=collection_name,
            embedding_function=embeddings
        )
        
        # 벡터 스토어를 전역에 저장
        st.session_state[global_vector_store_key] = vector_store
        
        logger.info("=== ChromaDB 벡터 스토어 로드 완료 ===")
        return vector_store
        
    except Exception as e:
        error_msg = f"ChromaDB 로드 실패: {e}"
        logger.error(f"ChromaDB 로드 실패: {e}", exc_info=True)
        st.error(f"❌ {error_msg}")
        return None

############################### 2단계 : RAG 기능 구현과 관련된 함수들 ##########################

## 하이브리드 서치 구현
class HybridRetriever:
    """BM25와 벡터 검색을 결합한 하이브리드 검색기"""
    
    def __init__(self, vector_store, documents: List[Document], k: int = 8, 
                 bm25_weight: float = 0.3, vector_weight: float = 0.7):
        self.vector_store = vector_store
        self.documents = documents
        self.k = k
        
        # BM25 검색기 초기화
        self.bm25_retriever = BM25Retriever.from_documents(documents)
        self.bm25_retriever.k = k
        
        # 벡터 검색기 초기화
        self.vector_retriever = vector_store.as_retriever(search_kwargs={"k": k})
        
        # 가중치 설정
        self.bm25_weight = bm25_weight
        self.vector_weight = vector_weight
        
        # 앙상블 검색기 생성
        self._create_ensemble_retriever()
    
    def _create_ensemble_retriever(self):
        """앙상블 검색기 생성"""
        self.ensemble_retriever = EnsembleRetriever(
            retrievers=[self.bm25_retriever, self.vector_retriever],
            weights=[self.bm25_weight, self.vector_weight]
        )
    
    def update_weights(self, bm25_weight: float, vector_weight: float):
        """가중치 업데이트"""
        self.bm25_weight = bm25_weight
        self.vector_weight = vector_weight
        self._create_ensemble_retriever()
        logger.info(f"가중치 업데이트 완료 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
    
    def search(self, query: str, k: int = 3, 
               bm25_weight: float = None, vector_weight: float = None) -> List[Document]:
        """하이브리드 검색 실행"""
        try:
            logger.info(f"=== 하이브리드 검색 시작: '{query[:50]}...' ===")
            
            # 전달받은 가중치가 있으면 업데이트
            if bm25_weight is not None and vector_weight is not None:
                self.update_weights(bm25_weight, vector_weight)
            
            logger.info(f"검색 가중치 - BM25: {self.bm25_weight:.2f}, 벡터: {self.vector_weight:.2f}")
            
            # 1차 검색: k=8개 결과
            logger.info("1차 검색 실행 중... (BM25 + 벡터 앙상블)")
            initial_results = self.ensemble_retriever.get_relevant_documents(query)
            logger.info(f"1차 검색 완료: {len(initial_results)}개 문서 발견")
            
            # 검색된 문서 정보 로깅
            for i, doc in enumerate(initial_results[:3]):  # 상위 3개만 로깅
                source = doc.metadata.get('source', 'Unknown')
                title = doc.metadata.get('title', 'No Title')
                logger.info(f"  문서 {i+1}: {source} | {title}")
            
            # 2차 재순위화: 관련성 점수 계산 및 상위 k개 선택
            logger.info("2차 재순위화(Reranking) 실행 중...")
            reranked_results = self._rerank_results(query, initial_results, k)
            
            logger.info(f"=== 하이브리드 검색 완료 ===")
            logger.info(f"  초기 결과: {len(initial_results)}개 → 최종 결과: {len(reranked_results)}개")
            
            # 최종 결과 상세 로깅
            for i, doc in enumerate(reranked_results):
                source = doc.metadata.get('source', 'Unknown')
                title = doc.metadata.get('title', 'No Title')
                relevance = doc.metadata.get('relevance_score', 'N/A')
                logger.info(f"  최종 {i+1}위: {source} | {title} | 점수: {relevance}")
            
            return reranked_results
            
        except Exception as e:
            logger.error(f"하이브리드 검색 실패: {e}")
            # 폴백: 벡터 검색만 사용
            logger.info("벡터 검색 폴백 실행...")
            try:
                logger.info("ChromaDB 벡터 검색 실행 중...")
                vector_results = self.vector_retriever.get_relevant_documents(query)
                logger.info(f"벡터 검색 폴백 완료: {len(vector_results)}개 문서")
                
                # 벡터 검색 결과 로깅
                for i, doc in enumerate(vector_results[:3]):
                    source = doc.metadata.get('source', 'Unknown')
                    title = doc.metadata.get('title', 'No Title')
                    logger.info(f"  벡터 검색 결과 {i+1}: {source} | {title}")
                
                return vector_results
            except Exception as vector_error:
                logger.error(f"벡터 검색 폴백도 실패: {vector_error}")
                return []
    
    def _rerank_results(self, query: str, documents: List[Document], k: int) -> List[Document]:
        """검색 결과 재순위화"""
        try:
            logger.info(f"재순위화 시작: {len(documents)}개 문서 처리 중...")
            scored_docs = []
            
            for i, doc in enumerate(documents):
                # 관련성 점수 계산 (간단한 키워드 매칭 + 메타데이터 가중치)
                score = self._calculate_relevance_score(query, doc)
                scored_docs.append((doc, score))
                
                # 상위 5개 문서의 점수만 로깅
                if i < 5:
                    source = doc.metadata.get('source', 'Unknown')
                    title = doc.metadata.get('title', 'No Title')
                    logger.info(f"  문서 {i+1} 점수: {score:.3f} | {source} | {title}")
            
            # 점수 기준 내림차순 정렬
            scored_docs.sort(key=lambda x: x[1], reverse=True)
            logger.info(f"재순위화 완료: 상위 {k}개 문서 선택")
            
            # 상위 k개 반환하고 관련성 점수를 메타데이터에 저장
            final_docs = []
            for doc, score in scored_docs[:k]:
                doc.metadata['relevance_score'] = round(score, 3)
                final_docs.append(doc)
            
            return final_docs
            
        except Exception as e:
            logger.warning(f"재순위화 실패: {e}")
            return documents[:k]
    
    def _calculate_relevance_score(self, query: str, doc: Document) -> float:
        """문서 관련성 점수 계산"""
        score = 0.0
        score_details = []
        
        try:
            # 1. 키워드 매칭 점수 (BM25 스타일)
            query_words = set(query.lower().split())
            doc_words = set(doc.page_content.lower().split())
            
            # 공통 단어 수
            common_words = query_words.intersection(doc_words)
            if common_words:
                keyword_score = len(common_words) * 0.1
                score += keyword_score
                score_details.append(f"키워드: +{keyword_score:.3f} ({len(common_words)}개 공통)")
            
            # 2. 메타데이터 가중치
            metadata = doc.metadata
            if metadata:
                # 최신 문서 우선
                if 'updated' in metadata:
                    try:
                        from datetime import datetime
                        updated_date = datetime.fromisoformat(metadata['updated'].replace('Z', '+00:00'))
                        days_old = (datetime.now().replace(tzinfo=updated_date.tzinfo) - updated_date).days
                        if days_old <= 30:  # 30일 이내
                            score += 0.2
                            score_details.append("최신성(30일): +0.200")
                        elif days_old <= 90:  # 90일 이내
                            score += 0.1
                            score_details.append("최신성(90일): +0.100")
                    except:
                        pass
                
                # 문서 타입별 가중치
                doctype_weights = {
                    'api': 0.3,
                    'spec': 0.3,
                    'guide': 0.2,
                    'blog': 0.1
                }
                
                if 'doctype' in metadata:
                    doc_type = metadata['doctype'].lower()
                    for dt, weight in doctype_weights.items():
                        if dt in doc_type:
                            score += weight
                            score_details.append(f"문서타입({dt}): +{weight:.3f}")
                            break
                
                # 파일 확장자별 가중치
                if 'source' in metadata:
                    source = metadata['source'].lower()
                    if source.endswith('.pdf'):
                        score += 0.1  # PDF 우선
                        score_details.append("파일형식(PDF): +0.100")
                    elif source.endswith(('.md', '.txt')):
                        score += 0.05
                        score_details.append("파일형식(TXT/MD): +0.050")
            
            # 3. 내용 길이 가중치 (적당한 길이 우선)
            content_length = len(doc.page_content)
            if 100 <= content_length <= 1000:
                score += 0.1
                score_details.append("내용길이(적당): +0.100")
            elif content_length > 1000:
                score += 0.05
                score_details.append("내용길이(긴): +0.050")
            
            # 점수 상세 정보 로깅 (디버그용)
            if score_details:
                source = metadata.get('source', 'Unknown') if metadata else 'Unknown'
                title = metadata.get('title', 'No Title') if metadata else 'No Title'
                logger.debug(f"점수 계산 상세 - {source} | {title}: {' + '.join(score_details)} = {score:.3f}")
            
        except Exception as e:
            logger.warning(f"점수 계산 중 오류: {e}")
        
        return score

def get_hybrid_retriever(vector_store, k: int = 8, 
                        bm25_weight: float = None, vector_weight: float = None) -> HybridRetriever:
    """하이브리드 검색기 생성"""
    try:
        logger.info("=== 하이브리드 검색기 생성 시작 ===")
        
        # 가중치 설정 (전달받은 값이 없으면 기본값 사용)
        if bm25_weight is None or vector_weight is None:
            config = get_hybrid_search_config()
            bm25_weight = config['bm25_weight']
            vector_weight = config['vector_weight']
            logger.info(f"기본 가중치 사용 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
        else:
            logger.info(f"전달받은 가중치 사용 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
        
        # 벡터 스토어에서 모든 문서 가져오기
        all_docs = []
        try:
            logger.info("ChromaDB에서 모든 문서 조회 중...")
            # ChromaDB에서 모든 문서 조회
            collection = vector_store._collection
            logger.info(f"ChromaDB 컬렉션 접근: {collection.name}")
            
            # 컬렉션 정보 확인
            total_count = collection.count()
            logger.info(f"컬렉션 총 문서 수: {total_count}")
            
            if total_count == 0:
                logger.warning("컬렉션에 문서가 없습니다")
                return None
            
            # 모든 문서 조회
            results = collection.get(include=['documents', 'metadatas'])
            logger.info(f"문서 조회 결과: {len(results['documents'])}개 문서")
            
            # 문서 메타데이터 샘플 로깅
            if results['metadatas']:
                sample_metadata = results['metadatas'][0]
                logger.info(f"메타데이터 샘플: {sample_metadata}")
            
            for i, (doc_content, metadata) in enumerate(zip(results['documents'], results['metadatas'])):
                doc = Document(
                    page_content=doc_content,
                    metadata=metadata
                )
                all_docs.append(doc)
                
                # 처음 3개 문서의 내용 미리보기 로깅
                if i < 3:
                    content_preview = doc_content[:100] + "..." if len(doc_content) > 100 else doc_content
                    source = metadata.get('source', 'Unknown')
                    title = metadata.get('title', 'No Title')
                    logger.info(f"문서 {i+1} 미리보기: {source} | {title} | 내용: {content_preview}")
                
        except Exception as e:
            logger.error(f"ChromaDB 문서 조회 실패: {e}", exc_info=True)
            all_docs = []
        
        logger.info(f"총 {len(all_docs)}개 문서 로드 완료")
        
        # 하이브리드 검색기 생성 (가중치 포함)
        logger.info("하이브리드 검색기 객체 생성 중...")
        hybrid_retriever = HybridRetriever(
            vector_store, all_docs, k, 
            bm25_weight=bm25_weight, 
            vector_weight=vector_weight
        )
        logger.info(f"=== 하이브리드 검색기 생성 완료: {len(all_docs)}개 문서, 가중치: BM25={bm25_weight:.2f}, Vector={vector_weight:.2f} ===")
        
        return hybrid_retriever
        
    except Exception as e:
        logger.error(f"하이브리드 검색기 생성 실패: {e}", exc_info=True)
        return None

# 전역 하이브리드 검색 설정 변수 추가 (파일 상단에)
HYBRID_SEARCH_CONFIG = {
    'bm25_weight': 0.3,
    'vector_weight': 0.7,
    'initial_k': 8,
    'final_k': 3,
    'enable_reranking': True,
    'metadata_boost': True,
    'recency_boost': True
}

def get_hybrid_search_config():
    """하이브리드 검색 설정 반환"""
    global HYBRID_SEARCH_CONFIG
    
    # UI에서 설정한 하이브리드 검색 설정 사용
    if 'hybrid_search_config' in st.session_state:
        config = st.session_state.hybrid_search_config
        # 전역 설정 업데이트
        HYBRID_SEARCH_CONFIG.update({
            'bm25_weight': config.get('bm25_weight', 0.3),
            'vector_weight': config.get('vector_weight', 0.7)
        })
        logger.info(f"세션 상태에서 가중치 업데이트 - BM25: {HYBRID_SEARCH_CONFIG['bm25_weight']:.2f}, Vector: {HYBRID_SEARCH_CONFIG['vector_weight']:.2f}")
    
    return HYBRID_SEARCH_CONFIG.copy()

def update_hybrid_search_config(config: dict):
    """하이브리드 검색 설정 업데이트"""
    global HYBRID_SEARCH_CONFIG
    
    try:
        # 전역 설정 업데이트
        HYBRID_SEARCH_CONFIG.update({
            'bm25_weight': config['bm25_weight'],
            'vector_weight': config['vector_weight'],
            'initial_k': config.get('initial_k', 8),
            'final_k': config.get('final_k', 3),
            'enable_reranking': config.get('enable_reranking', True),
            'metadata_boost': config.get('metadata_boost', True),
            'recency_boost': config.get('recency_boost', True)
        })
        
        # 세션 상태도 업데이트
        if 'hybrid_search_config' not in st.session_state:
            st.session_state['hybrid_search_config'] = {}
        
        st.session_state['hybrid_search_config'].update({
            'bm25_weight': config['bm25_weight'],
            'vector_weight': config['vector_weight'],
            'initial_k': config.get('initial_k', 8),
            'final_k': config.get('final_k', 3),
            'enable_reranking': config.get('enable_reranking', True),
            'metadata_boost': config.get('metadata_boost', True),
            'recency_boost': config.get('recency_boost', True)
        })
        
        logger.info(f"하이브리드 검색 설정 업데이트 완료 - BM25: {config['bm25_weight']:.2f}, Vector: {config['vector_weight']:.2f}")
        return True
    except Exception as e:
        logger.error(f"하이브리드 검색 설정 업데이트 실패: {e}")
        return False

def test_hybrid_search(query: str, vector_store, bm25_weight: float = None, vector_weight: float = None) -> dict:
    """하이브리드 검색 테스트 및 성능 측정"""
    try:
        import time
        
        # 가중치 설정 (전달받은 값이 없으면 현재 설정 사용)
        if bm25_weight is None or vector_weight is None:
            config = get_hybrid_search_config()
            bm25_weight = config['bm25_weight']
            vector_weight = config['vector_weight']
            logger.info(f"테스트에 현재 설정된 가중치 사용 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
        else:
            logger.info(f"테스트에 전달받은 가중치 사용 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
        
        # 벡터 검색만 테스트
        start_time = time.time()
        vector_retriever = vector_store.as_retriever(search_kwargs={"k": 3})
        vector_results = vector_retriever.invoke(query)
        vector_time = time.time() - start_time
        
        # 하이브리드 검색 테스트 (가중치 포함)
        start_time = time.time()
        hybrid_retriever = get_hybrid_retriever(
            vector_store, k=8, 
            bm25_weight=bm25_weight, 
            vector_weight=vector_weight
        )
        hybrid_results = hybrid_retriever.search(
            query, k=3,
            bm25_weight=bm25_weight,
            vector_weight=vector_weight
        )
        hybrid_time = time.time() - start_time
        
        # 성능 지표 계산
        performance_metrics = {
            'search_time': round(hybrid_time, 3),
            'result_count': len(hybrid_results),
            'avg_relevance': 'N/A',  # 관련성 점수는 메타데이터에 저장되어 있음
            'bm25_weight': bm25_weight,
            'vector_weight': vector_weight
        }
        
        # 결과에 관련성 점수 추가 (메타데이터에 저장된 경우)
        for doc in vector_results:
            if 'relevance_score' not in doc.metadata:
                doc.metadata['relevance_score'] = 'N/A'
        
        for doc in hybrid_results:
            if 'relevance_score' not in doc.metadata:
                doc.metadata['relevance_score'] = 'N/A'
        
        return {
            'vector_results': vector_results,
            'hybrid_results': hybrid_results,
            'performance_metrics': performance_metrics,
            'comparison': {
                'vector_time': round(vector_time, 3),
                'hybrid_time': round(hybrid_time, 3),
                'time_improvement': round((vector_time - hybrid_time) / vector_time * 100, 1) if vector_time > 0 else 0,
                'result_overlap': len(set([doc.metadata.get('file_name', '') for doc in vector_results]) & 
                                   set([doc.metadata.get('file_name', '') for doc in hybrid_results])),
                'weights_used': f"BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}"
            }
        }
        
    except Exception as e:
        logger.error(f"하이브리드 검색 테스트 실패: {e}")
        return None

## 사용자 질문에 대한 RAG 처리
def process_question(user_question, rag_chain=None):
    """사용자 질문을 처리하고 답변을 반환"""
    logger.info(f"=== 질문 처리 시작: {user_question[:50]}... ===")
    try:
        # RAG 체인이 전달되지 않았으면 가져오기
        if rag_chain is None:
            rag_chain = get_cached_rag_chain()
            if rag_chain is None:
                logger.error("RAG 체인 생성 실패")
                st.error("RAG 체인 생성에 실패했습니다.")
                return None, []
        
        # 질문 처리
        logger.info("=== RAG 체인 실행 시작 ===")
        response = rag_chain.invoke(user_question)
        logger.info(f"RAG 체인 실행 완료, 응답 길이: {len(response) if response else 0}")
        
        if response:
            logger.info(f"생성된 응답 미리보기: {response[:200]}...")
        
        # 관련 문서 검색 (하이브리드 검색 사용)
        logger.info("=== 관련 문서 검색 시작 ===")
        retrieve_docs = []
        try:
            logger.info("벡터 스토어 가져오기 시작...")
            vector_store = get_cached_vector_store()
            if vector_store:
                logger.info("벡터 스토어 로드 성공, 하이브리드 검색기 생성 시작...")
                
                # 현재 설정된 가중치 가져오기
                current_config = get_hybrid_search_config()
                bm25_weight = current_config['bm25_weight']
                vector_weight = current_config['vector_weight']
                logger.info(f"현재 설정된 가중치 - BM25: {bm25_weight:.2f}, Vector: {vector_weight:.2f}")
                
                # 하이브리드 검색기 사용 (현재 가중치로 생성)
                hybrid_retriever = get_hybrid_retriever(
                    vector_store, k=8, 
                    bm25_weight=bm25_weight, 
                    vector_weight=vector_weight
                )
                if hybrid_retriever:
                    logger.info("하이브리드 검색기 생성 성공, 검색 실행 중...")
                    # 검색 시에도 현재 가중치 전달
                    retrieve_docs = hybrid_retriever.search(
                        user_question, k=3,
                        bm25_weight=bm25_weight,
                        vector_weight=vector_weight
                    )
                    logger.info(f"하이브리드 검색 완료, 문서 수: {len(retrieve_docs)}")
                    
                    # 검색된 문서 상세 정보 로깅
                    for i, doc in enumerate(retrieve_docs):
                        source = doc.metadata.get('source', 'Unknown')
                        title = doc.metadata.get('title', 'No Title')
                        relevance = doc.metadata.get('relevance_score', 'N/A')
                        logger.info(f"  검색된 문서 {i+1}: {source} | {title} | 관련성: {relevance}")
                else:
                    # 폴백: 기존 벡터 검색
                    logger.info("하이브리드 검색기 생성 실패, 벡터 검색 폴백 실행...")
                    logger.info("ChromaDB 직접 검색 실행 중...")
                    retriever = vector_store.as_retriever(search_kwargs={"k": 3})
                    retrieve_docs = retriever.invoke(user_question)
                    logger.info(f"벡터 검색 폴백 완료, 문서 수: {len(retrieve_docs)}")
                    
                    # 폴백 검색 결과 로깅
                    for i, doc in enumerate(retrieve_docs):
                        source = doc.metadata.get('source', 'Unknown')
                        title = doc.metadata.get('title', 'No Title')
                        logger.info(f"  폴백 검색 문서 {i+1}: {source} | {title}")
            else:
                logger.warning("벡터 스토어를 가져올 수 없습니다")
        except Exception as e:
            logger.error(f"관련 문서 검색 실패: {e}", exc_info=True)

        logger.info("=== 질문 처리 완료 ===")
        logger.info(f"최종 결과: 응답 길이 {len(response) if response else 0}자, 문서 {len(retrieve_docs)}개")
        return response, retrieve_docs
    except Exception as e:
        logger.error(f"질문 처리 중 오류 발생: {str(e)}", exc_info=True)
        st.error(f"질문 처리 중 오류 발생: {str(e)}")
        return None, []

@st.cache_resource
def get_cached_vector_store():
    """캐시된 벡터 스토어 반환 (모델별 캐시)"""
    try:
        logger.info("=== 캐시된 벡터 스토어 확인 시작 ===")
        
        # 현재 선택된 모델과 임베딩 모델로 캐시 키 생성
        selected_model = st.session_state.get('selected_model', 'exaone3.5')
        selected_embedding = st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask')
        
        cache_key = f"vector_store_{selected_model}_{selected_embedding}"
        logger.info(f"캐시 키: {cache_key}")
        
        # 기존 캐시된 벡터 스토어가 있으면 재사용
        if cache_key in st.session_state:
            try:
                logger.info("기존 캐시된 벡터 스토어 확인 중...")
                vector_store = st.session_state[cache_key]
                # 간단한 테스트로 연결 상태 확인
                test_collection = vector_store._collection
                doc_count = test_collection.count()
                logger.info(f"기존 벡터 스토어 캐시 재사용 성공: {cache_key} (문서 수: {doc_count})")
                return vector_store
            except Exception as e:
                logger.warning(f"기존 벡터 스토어 캐시 재사용 실패: {e}")
                # 기존 캐시 제거
                del st.session_state[cache_key]
                logger.info("실패한 캐시 제거 완료")
        
        # 새 벡터 스토어 생성
        logger.info("새 벡터 스토어 생성 시작...")
        vector_store = load_chroma_store()
        if vector_store:
            st.session_state[cache_key] = vector_store
            logger.info(f"새 벡터 스토어 캐시 생성 완료: {cache_key}")
        else:
            logger.error("벡터 스토어 생성 실패")
        
        logger.info("=== 캐시된 벡터 스토어 확인 완료 ===")
        return vector_store
        
    except Exception as e:
        logger.error(f"벡터 스토어 캐시 생성 실패: {str(e)}", exc_info=True)
        return None

@st.cache_resource
def get_cached_rag_chain():
    """캐시된 RAG 체인 반환 (모델별 캐시)"""
    try:
        # 현재 선택된 모델로 캐시 키 생성
        selected_model = st.session_state.get('selected_model', 'exaone3.5')
        cache_key = f"rag_chain_{selected_model}"
        
        # 기존 캐시된 RAG 체인이 있으면 재사용
        if cache_key in st.session_state:
            logger.info(f"기존 RAG 체인 캐시 재사용: {cache_key}")
            return st.session_state[cache_key]
        
        # 새 RAG 체인 생성
        rag_chain = get_rag_chain()
        if rag_chain:
            st.session_state[cache_key] = rag_chain
            logger.info(f"새 RAG 체인 캐시 생성: {cache_key}")
        
        return rag_chain
        
    except Exception as e:
        logger.error(f"RAG 체인 캐시 생성 실패: {str(e)}", exc_info=True)
        return None

def get_rag_chain() -> Runnable:
    """RAG 체인 생성 (실제 구현)"""
    try:
        # 선택된 모델 가져오기
        selected_model = st.session_state.get('selected_model', 'exaone3.5')
        
        # Ollama LLM 초기화
        from langchain_ollama import OllamaLLM
        llm = OllamaLLM(
            model=selected_model,
            temperature=0.1,
            top_p=0.9,
            max_tokens=2048
        )
        
        # 캐시된 임베딩 모델 사용
        embeddings = get_embedding_model()
        if embeddings is None:
            logger.error("임베딩 모델 로드 실패")
            return None
        
        # 캐시된 벡터 스토어 사용
        vector_store = get_cached_vector_store()
        if vector_store is None:
            logger.error("벡터 스토어 로드 실패")
            return None
        
        # 프롬프트 템플릿
        from langchain_core.prompts import PromptTemplate
        template = """당신은 bizMOB Platform 전문가입니다. 
다음 컨텍스트를 사용하여 질문에 답변해주세요. 답변은 한글로 해주세요:

컨텍스트:
{context}

질문: {question}

답변:"""
        
        prompt = PromptTemplate(
            input_variables=["context", "question"],
            template=template
        )
        
        # RAG 체인 생성 (하이브리드 검색 사용)
        from langchain_core.runnables import RunnablePassthrough
        from langchain_core.output_parsers import StrOutputParser
        
        # 하이브리드 검색기 생성
        hybrid_retriever = get_hybrid_retriever(vector_store, k=8)
        if hybrid_retriever:
            # 하이브리드 검색을 사용하는 커스텀 검색기 (Runnable 인터페이스 구현)
            from langchain_core.runnables import Runnable
            from typing import Any, List
            
            class HybridSearchRetriever(Runnable):
                def __init__(self, hybrid_retriever):
                    super().__init__()
                    self.hybrid_retriever = hybrid_retriever
                
                def invoke(self, input_data: Any, config: Any = None) -> List[Document]:
                    """Runnable 인터페이스 구현"""
                    if isinstance(input_data, str):
                        query = input_data
                    elif isinstance(input_data, dict) and 'question' in input_data:
                        query = input_data['question']
                    else:
                        query = str(input_data)
                    
                    return self.hybrid_retriever.search(query, k=3)
                
                def get_relevant_documents(self, query: str) -> List[Document]:
                    """기존 retriever 인터페이스 호환성"""
                    return self.hybrid_retriever.search(query, k=3)
            
            search_retriever = HybridSearchRetriever(hybrid_retriever)
        else:
            # 폴백: 기존 벡터 검색
            search_retriever = vector_store.as_retriever(search_kwargs={"k": 3})
        
        chain = (
            {"context": search_retriever, "question": RunnablePassthrough()}
            | prompt
            | llm
            | StrOutputParser()
        )
        
        return chain
        
    except Exception as e:
        logger.error(f"RAG 체인 생성 중 오류: {str(e)}", exc_info=True)
        st.error(f"RAG 체인 생성 중 오류: {str(e)}")
        return None

############################### 3단계 : 응답결과와 문서를 함께 보도록 도와주는 함수 ##########################

@st.cache_data(show_spinner=False)
def convert_pdf_to_images(pdf_path: str, dpi: int = 250) -> List[str]:
    doc = fitz.open(pdf_path)  # 문서 열기
    image_paths = []
    
    # 이미지 저장용 폴더 생성
    output_folder = "PDF_이미지_bizmob"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for page_num in range(len(doc)):  # 각 페이지를 순회
        page = doc.load_page(page_num)  # 페이지 로드

        zoom = dpi / 72  # 72이 디폴트 DPI
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat) # type: ignore

        image_path = os.path.join(output_folder, f"page_{page_num + 1}.png")  # 페이지 이미지 저장 page_1.png, page_2.png, etc.
        pix.save(image_path)  # PNG 형태로 저장
        image_paths.append(image_path)  # 경로를 저장
        
    return image_paths

def display_pdf_page(image_path: str, page_number: int) -> None:
    try:
        image_bytes = open(image_path, "rb").read()  # 파일에서 이미지 인식
        st.image(image_bytes, caption=f"Page {page_number}", output_format="PNG", width=600)
    except Exception as e:
        st.error(f"이미지 표시 중 오류: {str(e)}")

def natural_sort_key(s):
    return [int(text) if text.isdigit() else text for text in re.split(r'(\d+)', s)]

def check_vector_db_exists():
    """벡터DB가 존재하는지 확인"""
    return os.path.exists(get_chroma_db_path())

def load_saved_model_info():
    """저장된 모델 정보를 불러옴"""
    try:
        if os.path.exists('vector_db_model_info.json'):
            import json
            with open('vector_db_model_info.json', 'r', encoding='utf-8') as f:
                model_info = json.load(f)
            return model_info
        return None
    except Exception as e:
        st.warning(f"⚠️ 저장된 모델 정보 불러오기 실패: {str(e)}")
        return None

def check_ollama_models():
    """사용 가능한 Ollama 모델 확인"""
    try:
        import subprocess
        result = subprocess.run(['ollama', 'list'], capture_output=True, text=True)
        if result.returncode == 0:
            return True
        return False
    except:
        return False

def get_ollama_models():
    """Ollama에서 사용 가능한 모델 목록을 가져옴"""
    try:
        import subprocess
        result = subprocess.run(['ollama', 'list'], capture_output=True, text=True)
        if result.returncode == 0:
            lines = result.stdout.strip().split('\n')
            models = []
            for line in lines[1:]:  # 첫 번째 줄은 헤더이므로 제외
                if line.strip():
                    parts = line.split()
                    if len(parts) >= 2:
                        model_name = parts[0]
                        model_size = parts[1] if len(parts) > 1 else "Unknown"
                        models.append({
                            'name': model_name,
                            'size': model_size
                        })
            return models
        return []
    except Exception as e:
        st.error(f"Ollama 모델 목록 가져오기 실패: {str(e)}")
        return []

def get_available_embedding_models():
    """사용 가능한 임베딩 모델 목록을 반환"""
    return {
        'sentence-transformers/all-MiniLM-L6-v2': {
            'name': 'all-MiniLM-L6-v2',
            'description': '영어 전용, 빠르고 가벼운 모델',
            'language': 'English',
            'size': 'Small'
        },
        'jhgan/ko-sroberta-multitask': {
            'name': 'ko-sroberta-multitask',
            'description': '한국어 전용, 다중 작업 모델',
            'language': 'Korean',
            'size': 'Medium'
        },
        'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2': {
            'name': 'paraphrase-multilingual-MiniLM-L12-v2',
            'description': '다국어 지원, 균형잡힌 성능',
            'language': 'Multilingual',
            'size': 'Medium'
        },
        'sentence-transformers/all-mpnet-base-v2': {
            'name': 'all-mpnet-base-v2',
            'description': '영어 전용, 고품질 임베딩',
            'language': 'English',
            'size': 'Large'
        },
        'intfloat/multilingual-e5-large': {
            'name': 'multilingual-e5-large',
            'description': '다국어 지원, 고품질 임베딩',
            'language': 'Multilingual',
            'size': 'Large'
        }
    }

class SafeSentenceTransformerEmbeddings(Embeddings):
    """safetensors를 직접 사용하는 안전한 임베딩 클래스"""
    
    def __init__(self, model_name: str, device: str = 'cpu'):
        self.model_name = model_name
        self.device = device
        self.model = None
        self.tokenizer = None
        self._load_model()
    
    def _load_model(self):
        """모델을 안전하게 로드 (safetensors 직접 사용)"""
        try:
            # 환경 변수 설정으로 safetensors 강제 사용
            import os
            os.environ['SAFETENSORS_FAST_GPU'] = '1'
            os.environ['TRANSFORMERS_OFFLINE'] = '0'
            os.environ['TORCH_WEIGHTS_ONLY'] = '1'
            os.environ['HF_HUB_DISABLE_TELEMETRY'] = '1'
            os.environ['TRANSFORMERS_USE_SAFETENSORS'] = '1'
            os.environ['TORCH_WARN_ON_LOAD'] = '0'
            os.environ['TORCH_LOAD_WARN_ONLY'] = '0'
            os.environ['TRANSFORMERS_SAFE_SERIALIZATION'] = '1'
            os.environ['HF_HUB_ENABLE_HF_TRANSFER'] = '1'
            
            # transformers와 safetensors를 직접 사용
            from transformers import AutoTokenizer, AutoModel
            import torch
            import safetensors
            
            # 토크나이저 로드
            self.tokenizer = AutoTokenizer.from_pretrained(
                self.model_name,
                trust_remote_code=True,
                use_safetensors=True
            )
            
            # 모델 로드 (safetensors 강제 사용)
            self.model = AutoModel.from_pretrained(
                self.model_name,
                trust_remote_code=True,
                use_safetensors=True,
                torch_dtype=torch.float32,
                low_cpu_mem_usage=True
            )
            
            self.model.to(self.device)
            self.model.eval()
            
        except Exception as e:
            st.error(f"모델 로딩 실패: {str(e)}")
            raise e
    
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """문서들을 임베딩"""
        if self.model is None:
            self._load_model()
        
        try:
            embeddings = []
            for text in texts:
                # 토크나이징
                inputs = self.tokenizer(
                    text,
                    return_tensors="pt",
                    max_length=512,
                    truncation=True,
                    padding=True
                )
                
                # GPU로 이동
                inputs = {k: v.to(self.device) for k, v in inputs.items()}
                
                # 임베딩 생성
                with torch.no_grad():
                    outputs = self.model(**inputs)
                    # 마지막 hidden state의 평균을 임베딩으로 사용
                    embedding = outputs.last_hidden_state.mean(dim=1)
                    embeddings.append(embedding.cpu().numpy().flatten().tolist())
            
            return embeddings
            
        except Exception as e:
            st.error(f"임베딩 생성 실패: {str(e)}")
            raise e
    
    def embed_query(self, text: str) -> List[float]:
        """단일 쿼리를 임베딩"""
        if self.model is None:
            self._load_model()
        
        try:
            # 토크나이징
            inputs = self.tokenizer(
                text,
                return_tensors="pt",
                max_length=512,
                truncation=True,
                padding=True
            )
            
            # GPU로 이동
            inputs = {k: v.to(self.device) for k, v in inputs.items()}
            
            # 임베딩 생성
            with torch.no_grad():
                outputs = self.model(**inputs)
                # 마지막 hidden state의 평균을 임베딩으로 사용
                embedding = outputs.last_hidden_state.mean(dim=1)
                return embedding.cpu().numpy().flatten().tolist()
                
        except Exception as e:
            st.error(f"쿼리 임베딩 실패: {str(e)}")
            raise e

@st.cache_resource
def get_embedding_model():
    """선택된 임베딩 모델을 반환 (safetensors 지원)"""
    selected_embedding = st.session_state.get('selected_embedding_model', 'jhgan/ko-sroberta-multitask')
    logger.info(f"=== 임베딩 모델 로드 시작: {selected_embedding} ===")
    
    try:
        logger.info("1. SafeSentenceTransformerEmbeddings 시도")
        logger.info(f"모델 이름: {selected_embedding}, 디바이스: cpu")
        
        # 커스텀 임베딩 클래스 사용
        embeddings = SafeSentenceTransformerEmbeddings(
            model_name=selected_embedding,
            device='cpu'
        )
        logger.info("SafeSentenceTransformerEmbeddings 로드 성공")
        logger.info(f"=== 임베딩 모델 로드 완료: {selected_embedding} ===")
        return embeddings
        
    except Exception as e:
        logger.warning(f"SafeSentenceTransformerEmbeddings 실패: {str(e)}")
        logger.info("HuggingFaceEmbeddings로 폴백 시도...")
        st.error(f"임베딩 모델 로딩 실패: {str(e)}")
        st.info("HuggingFaceEmbeddings로 재시도합니다...")
        
        try:
            logger.info("2. HuggingFaceEmbeddings 재시도")
            logger.info(f"모델 설정: device=cpu, torch_dtype=auto, low_cpu_mem_usage=True")
            
            # HuggingFaceEmbeddings로 fallback (safetensors 사용)
            embeddings = HuggingFaceEmbeddings(
                model_name=selected_embedding,
                model_kwargs={
                    'device': 'cpu',
                    'torch_dtype': 'auto',
                    'low_cpu_mem_usage': True,
                    'trust_remote_code': True
                },
                encode_kwargs={'normalize_embeddings': True}
            )
            
            logger.info("HuggingFaceEmbeddings 로드 성공")
            logger.info(f"=== 임베딩 모델 로드 완료 (폴백): {selected_embedding} ===")
            st.success(f"✅ {selected_embedding} 모델을 HuggingFaceEmbeddings로 로드했습니다.")
            return embeddings
            
        except Exception as e2:
            logger.error(f"HuggingFaceEmbeddings 재시도도 실패: {str(e2)}", exc_info=True)
            st.error(f"HuggingFaceEmbeddings 재시도도 실패: {str(e2)}")
            st.error("임베딩 모델을 로드할 수 없습니다. 다른 모델을 선택해주세요.")
            return None

def get_recommended_embedding_model(ai_model_name: str) -> str:
    """AI 모델에 따른 권장 임베딩 모델을 반환"""
    model_mapping = {
        'exaone3.5': 'jhgan/ko-sroberta-multitask',
        'llama3': 'sentence-transformers/all-mpnet-base-v2',
        'llama3.2': 'sentence-transformers/all-mpnet-base-v2',
        'llama3.2:3b': 'sentence-transformers/all-MiniLM-L6-v2',
        'llama3.2:8b': 'sentence-transformers/all-mpnet-base-v2',
        'llama3.2:70b': 'intfloat/multilingual-e5-large',
        'mistral': 'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2',
        'mistral:7b': 'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2',
        'mistral:instruct': 'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2',
        'codellama': 'sentence-transformers/all-mpnet-base-v2',
        'codellama:7b': 'sentence-transformers/all-MiniLM-L6-v2',
        'codellama:13b': 'sentence-transformers/all-mpnet-base-v2',
        'codellama:34b': 'intfloat/multilingual-e5-large',
        'gemma': 'sentence-transformers/all-mpnet-base-v2',
        'gemma:2b': 'sentence-transformers/all-MiniLM-L6-v2',
        'gemma:7b': 'sentence-transformers/all-mpnet-base-v2',
        'qwen': 'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2',
        'qwen:7b': 'sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2',
        'qwen:14b': 'intfloat/multilingual-e5-large',
        'phi': 'sentence-transformers/all-mpnet-base-v2',
        'phi:2.7b': 'sentence-transformers/all-MiniLM-L6-v2',
        'phi:3.5': 'sentence-transformers/all-mpnet-base-v2',
    }
    # 정확 매칭
    if ai_model_name in model_mapping:
        return model_mapping[ai_model_name]
    # 부분 매칭
    for key, value in model_mapping.items():
        if key in ai_model_name.lower():
            return value
    return 'jhgan/ko-sroberta-multitask'

def get_chroma_db_path():
    """ChromaDB 경로 반환"""
    chroma_path = "./chroma_db"
    logger.info(f"ChromaDB 경로 설정: {chroma_path}")
    return chroma_path

def get_model_info_path():
    """모델 정보 파일 경로 반환"""
    ai_model = st.session_state.get('selected_model', 'exaone3.5')
    import re
    safe_model = re.sub(r'[^a-zA-Z0-9_\-]', '_', ai_model)
    return f"vector_db_model_info_{safe_model}.json" 

def main():
    """메인 애플리케이션 (파일 관리 전용)"""
    logger.info("=== bizMOB 파일 관리 애플리케이션 시작 ===")
    
    st.set_page_config(
        page_title="bizMOB 파일 관리",
        page_icon="📁",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    logger.info("페이지 설정 완료")

    # 사이드바 - AI 모델 선택
    logger.info("AI 모델 선택 시작")
    if check_ollama_models():
        st.sidebar.success("✅ Ollama 연결됨")
        logger.info("Ollama 연결 성공")
        
        # 사용 가능한 모델 목록 가져오기
        available_models = get_ollama_models()
        logger.info(f"사용 가능한 모델 수: {len(available_models) if available_models else 0}")
        
        # 모델 선택기 표시
        show_model_selector(available_models, get_recommended_embedding_model, load_saved_model_info)
        logger.info("모델 선택기 표시 완료")
        
    else:
        st.sidebar.error("❌ Ollama 연결 실패")
        st.sidebar.info("Ollama가 설치되어 있고 실행 중인지 확인해주세요.")
        logger.warning("Ollama 연결 실패")
    
    # 사이드바 - 임베딩 모델 정보
    logger.info("임베딩 모델 정보 표시 시작")
    show_embedding_model_info(get_available_embedding_models, load_saved_model_info, get_recommended_embedding_model)
    logger.info("임베딩 모델 정보 표시 완료")
    
    # 채팅 페이지 링크
    st.sidebar.markdown("---")
    st.sidebar.markdown("### 💬 채팅 페이지")
    st.sidebar.markdown("[채팅 페이지로 이동](/Chat)")
    logger.info("채팅 페이지 링크 추가 완료")
    
    # 관리자 인터페이스 표시 (채팅 기능 제외)
    logger.info("관리자 인터페이스 표시 시작")
    show_admin_interface(
        None,  # display_chat_messages (None으로 설정하여 채팅 기능 비활성화)
        check_vector_db_exists, 
        initialize_vector_db_with_documents,
        None,  # add_chat_message (None으로 설정하여 채팅 기능 비활성화)
        None,  # process_question (None으로 설정하여 채팅 기능 비활성화)
        manage_uploaded_files, 
        load_saved_model_info
    )
    logger.info("관리자 인터페이스 표시 완료")
    
    logger.info("=== bizMOB 파일 관리 애플리케이션 로드 완료 ===")

if __name__ == "__main__":
    main() 