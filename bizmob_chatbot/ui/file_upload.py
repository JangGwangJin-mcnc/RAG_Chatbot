"""
파일 업로드 UI 컴포넌트
"""

import streamlit as st
import sys
import os
import logging
import glob
import pandas as pd
from typing import List

# 현재 디렉토리를 Python 경로에 추가
current_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(current_dir)

from utils.file_utils import save_uploaded_file
from core.vector_db_manager import VectorDBManager
from langchain_core.documents.base import Document
from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredExcelLoader, UnstructuredPowerPointLoader, UnstructuredWordDocumentLoader

# 로깅 설정
logger = logging.getLogger(__name__)


class FileUpload:
    """파일 업로드 관리"""
    
    def __init__(self):
        self.vector_db_manager = VectorDBManager()
    
    def display_file_upload_section(self):
        """파일 업로드 섹션 표시"""
        st.markdown("### 📁 파일 업로드")
        st.markdown("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
        
        # 파일 업로드 위젯
        uploaded_files = st.file_uploader(
            "문서 파일을 선택하세요",
            type=['pdf', 'xlsx', 'xls', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="여러 파일을 동시에 선택할 수 있습니다."
        )
        
        if uploaded_files:
            st.markdown("#### 📋 업로드된 파일 목록")
            
            success_count = 0
            error_count = 0
            
            for uploaded_file in uploaded_files:
                col1, col2, col3 = st.columns([3, 1, 1])
                
                with col1:
                    st.write(f"📄 {uploaded_file.name}")
                
                with col2:
                    # 파일 저장
                    file_path = save_uploaded_file(uploaded_file)
                    if file_path:
                        st.success("✅ 저장됨")
                        success_count += 1
                    else:
                        st.error("❌ 저장 실패")
                        error_count += 1
                
                with col3:
                    file_size = len(uploaded_file.getbuffer()) / 1024  # KB
                    st.write(f"{file_size:.1f} KB")
            
            # 업로드 결과 요약
            if success_count > 0:
                st.success(f"✅ {success_count}개 파일 업로드 완료")
                
                # 벡터 DB 자동 업데이트
                if st.button("🔄 벡터 DB 업데이트", type="primary"):
                    self.update_vector_db()
            else:
                st.error("❌ 파일 업로드에 실패했습니다")
    
    def load_all_documents_from_folder(self, folder_path: str = "PDF_bizMOB_Guide") -> List[Document]:
        """폴더에서 모든 문서 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        
        if not os.path.exists(folder_path):
            st.error(f"폴더가 존재하지 않습니다: {folder_path}")
            return documents
        
        # 지원하는 파일 확장자들
        supported_extensions = {
            '*.pdf': 'PDF',
            '*.xlsx': 'Excel',
            '*.xls': 'Excel',
            '*.pptx': 'PowerPoint',
            '*.ppt': 'PowerPoint',
            '*.docx': 'Word',
            '*.doc': 'Word'
        }
        
        all_files = []
        for pattern, file_type in supported_extensions.items():
            files = glob.glob(os.path.join(folder_path, pattern))
            all_files.extend([(f, file_type) for f in files])
        
        if not all_files:
            st.warning(f"{folder_path} 폴더에 지원하는 문서 파일이 없습니다.")
            st.info("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
            return documents
        
        for file_path, file_type in all_files:
            try:
                st.info(f"{file_type} 파일 로딩 중: {os.path.basename(file_path)}")
                
                if file_type == 'PDF':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 PDF 로드
                    loader = PyMuPDFLoader(file_path)
                    doc = loader.load()
                elif file_type == 'Excel':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 Excel 로드
                    loader = UnstructuredExcelLoader(file_path)
                    doc = loader.load()
                elif file_type == 'PowerPoint':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 PowerPoint 로드
                    loader = UnstructuredPowerPointLoader(file_path)
                    doc = loader.load()
                elif file_type == 'Word':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 Word 로드
                    loader = UnstructuredWordDocumentLoader(file_path)
                    doc = loader.load()
                
                # 메타데이터 추가
                for d in doc:
                    d.metadata['file_path'] = file_path
                    d.metadata['file_name'] = os.path.basename(file_path)
                    d.metadata['file_type'] = file_type
                
                documents.extend(doc)
                st.success(f"✅ {os.path.basename(file_path)} ({file_type}) 로딩 완료")
            except Exception as e:
                st.error(f"❌ {os.path.basename(file_path)} ({file_type}) 로딩 실패: {str(e)}")
        
        return documents
    
    def load_excel_file(self, file_path: str) -> List[Document]:
        """Excel 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # Excel 파일 읽기
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 데이터프레임을 텍스트로 변환
                text_content = f"시트명: {sheet_name}\n\n"
                
                # 헤더 정보 추가
                if not df.empty:
                    text_content += f"컬럼: {', '.join(df.columns.tolist())}\n\n"
                    
                    # 데이터 내용 추가 (처음 100행까지만)
                    for idx, row in df.head(100).iterrows():
                        row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                        if row_text.strip():
                            text_content += f"행 {idx+1}: {row_text}\n"
                
                # Document 생성
                doc = Document(
                    page_content=text_content,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_type': 'Excel',
                        'sheet_name': sheet_name
                    }
                )
                documents.append(doc)
                
        except Exception as e:
            st.error(f"Excel 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def load_powerpoint_file(self, file_path: str) -> List[Document]:
        """PowerPoint 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # PowerPoint 파일 읽기
            from pptx import Presentation
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = f"슬라이드 {slide_num}:\n\n"
                
                # 슬라이드의 모든 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += f"{shape.text}\n\n"
                
                # 텍스트가 있는 경우에만 Document 생성
                if text_content.strip() and text_content != f"슬라이드 {slide_num}:\n\n":
                    doc = Document(
                        page_content=text_content,
                        metadata={
                            'file_path': file_path,
                            'file_name': os.path.basename(file_path),
                            'file_type': 'PowerPoint',
                            'slide_number': slide_num
                        }
                    )
                    documents.append(doc)
                    
        except Exception as e:
            st.error(f"PowerPoint 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def load_word_file(self, file_path: str) -> List[Document]:
        """Word 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # Word 파일 읽기
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            
            # 전체 문서를 하나의 Document로 처리
            text_content = ""
            
            # 제목 정보 추가
            if doc.core_properties.title:
                text_content += f"제목: {doc.core_properties.title}\n\n"
            
            # 단락별로 텍스트 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content += para.text + "\n\n"
            
            # 표 내용 추출
            for table in doc.tables:
                text_content += "표 내용:\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                    if row_text.strip():
                        text_content += row_text + "\n"
                text_content += "\n"
            
            # 텍스트가 있는 경우에만 Document 생성
            if text_content.strip():
                doc_chunk = Document(
                    page_content=text_content,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_type': 'Word',
                        'title': doc.core_properties.title or 'Unknown',
                        'author': doc.core_properties.author or 'Unknown'
                    }
                )
                documents.append(doc_chunk)
                
        except Exception as e:
            st.error(f"Word 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def update_vector_db(self):
        """벡터 DB 업데이트"""
        try:
            with st.spinner("벡터 DB를 업데이트하는 중..."):
                # 진행 상황 표시
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # 1단계: 문서 로드
                status_text.text("📄 문서를 로드하는 중...")
                documents = self.load_all_documents_from_folder()
                progress_bar.progress(25)
                
                if not documents:
                    st.error("로드할 문서가 없습니다.")
                    return
                
                st.info(f"📄 {len(documents)}개의 문서를 로드했습니다")
                
                # 2단계: 문서 청킹
                status_text.text("✂️ 문서를 청킹하는 중...")
                progress_bar.progress(50)
                
                chunked_documents = self.vector_db_manager.chunk_documents(documents)
                st.info(f"✂️ {len(chunked_documents)}개의 청크로 분할 완료")
                
                # 3단계: 임베딩 모델 로드
                status_text.text("🤖 임베딩 모델을 로드하는 중...")
                progress_bar.progress(75)
                
                # 4단계: 벡터 DB 저장
                status_text.text("💾 벡터 DB에 저장하는 중...")
                success = self.vector_db_manager.save_to_vector_store(chunked_documents)
                progress_bar.progress(100)
                
                if success:
                    # 모델 정보 저장
                    self.vector_db_manager.save_model_info()
                    st.success("✅ 벡터 DB 업데이트가 완료되었습니다!")
                    logger.info("벡터 DB 업데이트 완료")
                else:
                    st.error("❌ 벡터 DB 업데이트에 실패했습니다")
                    logger.error("벡터 DB 업데이트 실패")
                    
        except Exception as e:
            st.error(f"벡터 DB 업데이트 중 오류가 발생했습니다: {str(e)}")
            logger.error(f"벡터 DB 업데이트 오류: {str(e)}")
    
    def display_uploaded_files_info(self):
        """업로드된 파일 정보 표시"""
        folder_path = "PDF_bizMOB_Guide"
        if not os.path.exists(folder_path):
            st.info("📁 업로드된 파일이 없습니다")
            return
        
        files = []
        total_size = 0
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if os.path.isfile(file_path):
                file_size = os.path.getsize(file_path)
                files.append({
                    'name': file,
                    'size': file_size
                })
                total_size += file_size
        
        if files:
            st.markdown("#### 📋 업로드된 파일 목록")
            st.info(f"📁 총 {len(files)}개 파일, {total_size / 1024:.1f} KB")
            
            for file_info in files:
                col1, col2 = st.columns([3, 1])
                with col1:
                    st.write(f"📄 {file_info['name']}")
                with col2:
                    st.write(f"{file_info['size'] / 1024:.1f} KB")
        else:
            st.info("📁 업로드된 파일이 없습니다") 