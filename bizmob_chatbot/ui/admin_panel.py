"""
관리자 패널 UI 컴포넌트
"""

import streamlit as st
import sys
import os
import logging
import glob
import pandas as pd
from typing import List

# 현재 디렉토리를 Python 경로에 추가
current_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.append(current_dir)

from core.auth import AuthManager
from core.vector_db_manager import VectorDBManager
from langchain_core.documents.base import Document
from langchain_community.document_loaders import PyMuPDFLoader, UnstructuredExcelLoader, UnstructuredPowerPointLoader, UnstructuredWordDocumentLoader

# 로깅 설정
logger = logging.getLogger(__name__)


class AdminPanel:
    """관리자 패널"""
    
    def __init__(self):
        self.auth_manager = AuthManager()
        self.vector_db_manager = VectorDBManager()
    
    def display_login_section(self):
        """로그인 섹션 표시"""
        st.sidebar.markdown("### 🔐 관리자 로그인")
        
        password = st.sidebar.text_input("관리자 비밀번호", type="password")
        if st.sidebar.button("로그인"):
            if self.auth_manager.login(password):
                st.sidebar.success("관리자로 로그인되었습니다!")
                st.rerun()
            else:
                st.sidebar.error("비밀번호가 올바르지 않습니다.")
    
    def require_admin_access(self):
        """관리자 접근 권한 확인"""
        if not self.auth_manager.is_admin():
            st.error("관리자 권한이 필요합니다.")
            st.stop()
    
    def display_admin_controls(self):
        """관리자 컨트롤 표시"""
        if not self.auth_manager.is_admin():
            return
        
        st.markdown("### 🔧 관리자 도구")
        
        # 탭 생성
        tab1, tab2, tab3 = st.tabs(["📊 벡터 DB 관리", "📁 파일 관리", "🔍 벡터 DB 데이터"])
        
        with tab1:
            self.display_vector_db_management()
        
        with tab2:
            self.display_file_management()
        
        with tab3:
            # This tab was moved to main.py's display_main_chat_interface
            # self.display_vector_db_data() 
            pass # Placeholder, as it's now handled in main.py
    
    def display_vector_db_management(self):
        """벡터 DB 관리 섹션"""
        st.markdown("#### 📊 벡터 DB 관리")
        
        # 벡터 DB 재생성 버튼
        if st.button("🔄 벡터 DB 재생성", type="primary"):
            self.rebuild_vector_db()
    
    def rebuild_vector_db(self):
        """벡터 DB 재생성"""
        try:
            with st.spinner("벡터 DB를 재생성하는 중..."):
                # 진행 상황 표시
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # 1단계: 모든 문서 로드
                status_text.text("📄 모든 문서를 로드하는 중...")
                progress_bar.progress(25)
                
                documents = self.load_all_documents_from_folder()
                
                if not documents:
                    st.error("로드할 문서가 없습니다.")
                    return
                
                st.info(f"📄 {len(documents)}개의 문서를 로드했습니다")
                
                # 2단계: 문서 청킹
                status_text.text("✂️ 문서를 청킹하는 중...")
                progress_bar.progress(50)
                
                chunked_documents = self.vector_db_manager.chunk_documents(documents)
                st.info(f"✂️ {len(chunked_documents)}개의 청크로 분할 완료")
                
                # 3단계: 벡터 DB 저장
                status_text.text("💾 벡터 DB에 저장하는 중...")
                progress_bar.progress(75)
                
                success = self.vector_db_manager.save_to_vector_store(chunked_documents)
                progress_bar.progress(100)
                
                if success:
                    # 모델 정보 저장
                    self.vector_db_manager.save_model_info()
                    st.success("✅ 벡터 DB 재생성이 완료되었습니다!")
                    logger.info("벡터 DB 재생성 완료")
                else:
                    st.error("❌ 벡터 DB 재생성에 실패했습니다")
                    logger.error("벡터 DB 재생성 실패")
                    
        except Exception as e:
            st.error(f"벡터 DB 재생성 중 오류가 발생했습니다: {str(e)}")
            logger.error(f"벡터 DB 재생성 오류: {str(e)}")
    
    def load_all_documents_from_folder(self, folder_path: str = "PDF_bizMOB_Guide") -> List[Document]:
        """폴더에서 모든 문서 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        
        if not os.path.exists(folder_path):
            st.error(f"폴더가 존재하지 않습니다: {folder_path}")
            return documents
        
        # 지원하는 파일 확장자들
        supported_extensions = {
            '*.pdf': 'PDF',
            '*.xlsx': 'Excel',
            '*.xls': 'Excel',
            '*.pptx': 'PowerPoint',
            '*.ppt': 'PowerPoint',
            '*.docx': 'Word',
            '*.doc': 'Word'
        }
        
        all_files = []
        for pattern, file_type in supported_extensions.items():
            files = glob.glob(os.path.join(folder_path, pattern))
            all_files.extend([(f, file_type) for f in files])
        
        if not all_files:
            st.warning(f"{folder_path} 폴더에 지원하는 문서 파일이 없습니다.")
            st.info("지원 형식: PDF, Excel (.xlsx, .xls), PowerPoint (.pptx, .ppt), Word (.docx, .doc)")
            return documents
        
        for file_path, file_type in all_files:
            try:
                st.info(f"{file_type} 파일 로딩 중: {os.path.basename(file_path)}")
                
                if file_type == 'PDF':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 PDF 로드
                    loader = PyMuPDFLoader(file_path)
                    doc = loader.load()
                elif file_type == 'Excel':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 Excel 로드
                    loader = UnstructuredExcelLoader(file_path)
                    doc = loader.load()
                elif file_type == 'PowerPoint':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 PowerPoint 로드
                    loader = UnstructuredPowerPointLoader(file_path)
                    doc = loader.load()
                elif file_type == 'Word':
                    # 기존 bizmob_chatbot.py와 동일한 방식으로 Word 로드
                    loader = UnstructuredWordDocumentLoader(file_path)
                    doc = loader.load()
                
                # 메타데이터 추가
                for d in doc:
                    d.metadata['file_path'] = file_path
                    d.metadata['file_name'] = os.path.basename(file_path)
                    d.metadata['file_type'] = file_type
                
                documents.extend(doc)
                st.success(f"✅ {os.path.basename(file_path)} ({file_type}) 로딩 완료")
            except Exception as e:
                st.error(f"❌ {os.path.basename(file_path)} ({file_type}) 로딩 실패: {str(e)}")
        
        return documents
    
    def load_excel_file(self, file_path: str) -> List[Document]:
        """Excel 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # Excel 파일 읽기
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # 데이터프레임을 텍스트로 변환
                text_content = f"시트명: {sheet_name}\n\n"
                
                # 헤더 정보 추가
                if not df.empty:
                    text_content += f"컬럼: {', '.join(df.columns.tolist())}\n\n"
                    
                    # 데이터 내용 추가 (처음 100행까지만)
                    for idx, row in df.head(100).iterrows():
                        row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                        if row_text.strip():
                            text_content += f"행 {idx+1}: {row_text}\n"
                
                # Document 생성
                doc = Document(
                    page_content=text_content,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_type': 'Excel',
                        'sheet_name': sheet_name
                    }
                )
                documents.append(doc)
                
        except Exception as e:
            st.error(f"Excel 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def load_powerpoint_file(self, file_path: str) -> List[Document]:
        """PowerPoint 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # PowerPoint 파일 읽기
            from pptx import Presentation
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text_content = f"슬라이드 {slide_num}:\n\n"
                
                # 슬라이드의 모든 텍스트 추출
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += f"{shape.text}\n\n"
                
                # 텍스트가 있는 경우에만 Document 생성
                if text_content.strip() and text_content != f"슬라이드 {slide_num}:\n\n":
                    doc = Document(
                        page_content=text_content,
                        metadata={
                            'file_path': file_path,
                            'file_name': os.path.basename(file_path),
                            'file_type': 'PowerPoint',
                            'slide_number': slide_num
                        }
                    )
                    documents.append(doc)
                    
        except Exception as e:
            st.error(f"PowerPoint 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def load_word_file(self, file_path: str) -> List[Document]:
        """Word 파일을 로드하여 Document 리스트로 변환"""
        documents = []
        try:
            # Word 파일 읽기
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            
            # 전체 문서를 하나의 Document로 처리
            text_content = ""
            
            # 제목 정보 추가
            if doc.core_properties.title:
                text_content += f"제목: {doc.core_properties.title}\n\n"
            
            # 단락별로 텍스트 추출
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content += para.text + "\n\n"
            
            # 표 내용 추출
            for table in doc.tables:
                text_content += "표 내용:\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                    if row_text.strip():
                        text_content += row_text + "\n"
                text_content += "\n"
            
            # 텍스트가 있는 경우에만 Document 생성
            if text_content.strip():
                doc_chunk = Document(
                    page_content=text_content,
                    metadata={
                        'file_path': file_path,
                        'file_name': os.path.basename(file_path),
                        'file_type': 'Word',
                        'title': doc.core_properties.title or 'Unknown',
                        'author': doc.core_properties.author or 'Unknown'
                    }
                )
                documents.append(doc_chunk)
                
        except Exception as e:
            st.error(f"Word 파일 처리 중 오류: {str(e)}")
        
        return documents
    
    def display_file_management(self):
        """파일 관리 섹션"""
        st.markdown("#### 📁 파일 관리")
        
        # 업로드된 파일 목록 표시
        uploaded_files = self.get_uploaded_files()
        
        if uploaded_files:
            st.info(f"📋 업로드된 파일: {len(uploaded_files)}개")
            
            # 파일 선택 기능
            st.markdown("#### 🔄 선택한 파일로 벡터 DB 업데이트")
            selected_files = st.multiselect(
                "벡터 DB에 추가할 파일을 선택하세요:",
                options=uploaded_files,
                help="여러 파일을 선택할 수 있습니다"
            )
            
            if selected_files:
                st.info(f"📄 선택된 파일: {len(selected_files)}개")
                for file_name in selected_files:
                    st.write(f"• {file_name}")
                
                if st.button("🔄 선택한 파일로 벡터 DB 업데이트", type="primary"):
                    self.update_vector_db_with_selected_files(selected_files)
            
            # 전체 파일 목록 표시
            st.markdown("#### 📋 전체 파일 목록")
            for file_name in uploaded_files:
                st.write(f"• {file_name}")
        else:
            st.warning("📁 업로드된 파일이 없습니다")
    
    def get_uploaded_files(self):
        """업로드된 파일 목록 반환"""
        folder_path = "PDF_bizMOB_Guide"
        if not os.path.exists(folder_path):
            return []
        
        files = []
        for file in os.listdir(folder_path):
            file_path = os.path.join(folder_path, file)
            if os.path.isfile(file_path):
                files.append(file)
        
        return files
    
    def update_vector_db_with_selected_files(self, selected_files):
        """선택한 파일로 벡터 DB 업데이트"""
        try:
            with st.spinner("선택한 파일로 벡터 DB를 업데이트하는 중..."):
                # 진행 상황 표시
                progress_bar = st.progress(0)
                status_text = st.empty()
                
                # 1단계: 선택한 파일에서 문서 로드
                status_text.text("📄 선택한 파일에서 문서를 로드하는 중...")
                progress_bar.progress(25)
                
                documents = []
                folder_path = "PDF_bizMOB_Guide"
                
                for file_name in selected_files:
                    file_path = os.path.join(folder_path, file_name)
                    if os.path.exists(file_path):
                        # 파일 타입에 따라 로드
                        file_ext = os.path.splitext(file_name)[1].lower()
                        
                        try:
                            if file_ext == '.pdf':
                                loader = PyMuPDFLoader(file_path)
                                file_docs = loader.load()
                            elif file_ext in ['.xlsx', '.xls']:
                                loader = UnstructuredExcelLoader(file_path)
                                file_docs = loader.load()
                            elif file_ext in ['.pptx', '.ppt']:
                                loader = UnstructuredPowerPointLoader(file_path)
                                file_docs = loader.load()
                            elif file_ext in ['.docx', '.doc']:
                                loader = UnstructuredWordDocumentLoader(file_path)
                                file_docs = loader.load()
                            else:
                                # 텍스트 파일로 처리
                                with open(file_path, 'r', encoding='utf-8') as f:
                                    content = f.read()
                                file_docs = [Document(page_content=content, metadata={'source': file_name})]
                            
                            documents.extend(file_docs)
                            st.success(f"✅ {file_name} 로드 완료")
                            
                        except Exception as e:
                            st.error(f"❌ {file_name} 로드 실패: {str(e)}")
                
                if not documents:
                    st.error("로드할 문서가 없습니다.")
                    return
                
                st.info(f"📄 {len(documents)}개의 문서를 로드했습니다")
                
                # 2단계: 문서 청킹
                status_text.text("✂️ 문서를 청킹하는 중...")
                progress_bar.progress(50)
                
                chunked_documents = self.vector_db_manager.chunk_documents(documents)
                st.info(f"✂️ {len(chunked_documents)}개의 청크로 분할 완료")
                
                # 3단계: 임베딩 모델 로드
                status_text.text("🤖 임베딩 모델을 로드하는 중...")
                progress_bar.progress(75)
                
                # 4단계: 벡터 DB 저장
                status_text.text("💾 벡터 DB에 저장하는 중...")
                success = self.vector_db_manager.save_to_vector_store(chunked_documents)
                progress_bar.progress(100)
                
                if success:
                    # 모델 정보 저장
                    self.vector_db_manager.save_model_info()
                    st.success("✅ 선택한 파일로 벡터 DB 업데이트가 완료되었습니다!")
                    logger.info("선택한 파일로 벡터 DB 업데이트 완료")
                else:
                    st.error("❌ 벡터 DB 업데이트에 실패했습니다")
                    logger.error("선택한 파일로 벡터 DB 업데이트 실패")
                    
        except Exception as e:
            st.error(f"벡터 DB 업데이트 중 오류가 발생했습니다: {str(e)}")
            logger.error(f"선택한 파일로 벡터 DB 업데이트 오류: {str(e)}") 