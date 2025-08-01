#!/usr/bin/env python3
"""
파일 처리 유틸리티
다양한 문서 형식을 처리하고 텍스트를 추출하는 기능
"""

import os
import io
import tempfile
from typing import List, Dict, Any
from langchain_core.documents import Document
import warnings

# 경고 억제
warnings.filterwarnings("ignore")

def get_supported_extensions() -> List[str]:
    """지원되는 파일 확장자 목록 반환"""
    return ['pdf', 'txt', 'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt']

def process_file(uploaded_file) -> List[Document]:
    """업로드된 파일을 처리하여 Document 객체 리스트 반환"""
    
    file_extension = uploaded_file.name.split('.')[-1].lower()
    
    if file_extension == 'pdf':
        return process_pdf(uploaded_file)
    elif file_extension in ['docx', 'doc']:
        return process_word(uploaded_file)
    elif file_extension in ['xlsx', 'xls']:
        return process_excel(uploaded_file)
    elif file_extension in ['pptx', 'ppt']:
        return process_powerpoint(uploaded_file)
    elif file_extension == 'txt':
        return process_text(uploaded_file)
    else:
        raise ValueError(f"지원되지 않는 파일 형식: {file_extension}")

def process_pdf(uploaded_file) -> List[Document]:
    """PDF 파일 처리"""
    try:
        import PyPDF2
        
        # 파일 내용 읽기
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
        
        documents = []
        for page_num, page in enumerate(pdf_reader.pages):
            text = page.extract_text()
            if text.strip():
                doc = Document(
                    page_content=text,
                    metadata={
                        'source': uploaded_file.name,
                        'page': page_num + 1,
                        'file_type': 'pdf'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"PDF 파일 처리 실패: {e}")

def process_word(uploaded_file) -> List[Document]:
    """Word 파일 처리"""
    try:
        from docx import Document as DocxDocument
        
        # 파일 내용 읽기
        doc = DocxDocument(io.BytesIO(uploaded_file.read()))
        
        documents = []
        text_content = []
        
        # 모든 단락에서 텍스트 추출
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # 표에서 텍스트 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(' | '.join(row_text))
        
        # 전체 텍스트를 하나의 문서로 생성
        if text_content:
            full_text = '\n\n'.join(text_content)
            doc = Document(
                page_content=full_text,
                metadata={
                    'source': uploaded_file.name,
                    'file_type': 'word'
                }
            )
            documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"Word 파일 처리 실패: {e}")

def process_excel(uploaded_file) -> List[Document]:
    """Excel 파일 처리"""
    try:
        import pandas as pd
        
        # 파일 내용 읽기
        excel_data = pd.read_excel(io.BytesIO(uploaded_file.read()), sheet_name=None)
        
        documents = []
        
        for sheet_name, df in excel_data.items():
            if not df.empty:
                # DataFrame을 텍스트로 변환
                text_content = []
                
                # 헤더 추가
                headers = df.columns.tolist()
                text_content.append(' | '.join(str(h) for h in headers))
                
                # 데이터 행 추가
                for _, row in df.iterrows():
                    row_text = ' | '.join(str(cell) for cell in row.values)
                    text_content.append(row_text)
                
                sheet_text = '\n'.join(text_content)
                
                doc = Document(
                    page_content=sheet_text,
                    metadata={
                        'source': uploaded_file.name,
                        'sheet': sheet_name,
                        'file_type': 'excel'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"Excel 파일 처리 실패: {e}")

def process_powerpoint(uploaded_file) -> List[Document]:
    """PowerPoint 파일 처리"""
    try:
        from pptx import Presentation
        
        # 파일 내용 읽기
        prs = Presentation(io.BytesIO(uploaded_file.read()))
        
        documents = []
        
        for slide_num, slide in enumerate(prs.slides):
            text_content = []
            
            # 슬라이드의 모든 도형에서 텍스트 추출
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
            
            # 텍스트가 있는 경우에만 문서 생성
            if text_content:
                slide_text = '\n\n'.join(text_content)
                doc = Document(
                    page_content=slide_text,
                    metadata={
                        'source': uploaded_file.name,
                        'slide': slide_num + 1,
                        'file_type': 'powerpoint'
                    }
                )
                documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"PowerPoint 파일 처리 실패: {e}")

def process_text(uploaded_file) -> List[Document]:
    """텍스트 파일 처리"""
    try:
        # 파일 내용 읽기
        content = uploaded_file.read().decode('utf-8')
        
        # 텍스트를 청크로 분할 (긴 텍스트의 경우)
        chunks = split_text_into_chunks(content, max_chunk_size=4000)
        
        documents = []
        for i, chunk in enumerate(chunks):
            doc = Document(
                page_content=chunk,
                metadata={
                    'source': uploaded_file.name,
                    'chunk': i + 1,
                    'file_type': 'text'
                }
            )
            documents.append(doc)
        
        return documents
        
    except Exception as e:
        raise Exception(f"텍스트 파일 처리 실패: {e}")

def split_text_into_chunks(text: str, max_chunk_size: int = 4000, overlap: int = 200) -> List[str]:
    """텍스트를 청크로 분할"""
    if len(text) <= max_chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + max_chunk_size
        
        # 문장 경계에서 분할 시도
        if end < len(text):
            # 마지막 문장 경계 찾기
            last_period = text.rfind('.', start, end)
            last_newline = text.rfind('\n', start, end)
            
            if last_period > start and last_period > last_newline:
                end = last_period + 1
            elif last_newline > start:
                end = last_newline + 1
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # 오버랩을 고려한 다음 시작점
        start = end - overlap
        if start >= len(text):
            break
    
    return chunks

def get_file_info(uploaded_file) -> Dict[str, Any]:
    """파일 정보 반환"""
    return {
        'name': uploaded_file.name,
        'size': uploaded_file.size,
        'type': uploaded_file.type,
        'extension': uploaded_file.name.split('.')[-1].lower() if '.' in uploaded_file.name else ''
    }

def validate_file(uploaded_file) -> bool:
    """파일 유효성 검사"""
    if uploaded_file is None:
        return False
    
    file_extension = uploaded_file.name.split('.')[-1].lower() if '.' in uploaded_file.name else ''
    supported_extensions = get_supported_extensions()
    
    return file_extension in supported_extensions

def process_multiple_files(uploaded_files) -> List[Document]:
    """여러 파일 처리"""
    all_documents = []
    
    for uploaded_file in uploaded_files:
        try:
            if validate_file(uploaded_file):
                documents = process_file(uploaded_file)
                all_documents.extend(documents)
            else:
                print(f"지원되지 않는 파일 형식: {uploaded_file.name}")
        except Exception as e:
            print(f"파일 처리 실패 {uploaded_file.name}: {e}")
    
    return all_documents

if __name__ == "__main__":
    # 테스트 코드
    print("파일 처리 유틸리티")
    print(f"지원되는 확장자: {get_supported_extensions()}") 